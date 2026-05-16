"""
Tread Customers PPTX Deck v3
- Bigger text (11pt body)
- Colored initial badge as visual "logo"
- Real company connections where known
- 4M: health → yellow, tickets added
- 2-column layout for readability
"""

import os, io, json, ssl, urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Color palette ──────────────────────────────────────────────────────────
HEALTH = {
    'green':  RGBColor(52, 168, 83),
    'yellow': RGBColor(230, 160, 0),
    'red':    RGBColor(210, 48, 37),
    'gray':   RGBColor(140, 140, 140),
}
TYPE_COLOR = {
    'Hauler':       RGBColor(0, 140, 186),
    'Producer':     RGBColor(118, 63, 166),
    'Construction': RGBColor(204, 88, 0),
    'Agriculture':  RGBColor(76, 153, 0),
    'Mixed':        RGBColor(60, 110, 130),
}
NAVY   = RGBColor(22, 48, 82)
WHITE  = RGBColor(255, 255, 255)
LGRAY  = RGBColor(246, 246, 246)
MGRAY  = RGBColor(208, 208, 208)
DGRAY  = RGBColor(80, 80, 80)
DARK   = RGBColor(28, 28, 28)
MED    = RGBColor(100, 100, 100)
LINK   = RGBColor(26, 115, 232)
WARN   = RGBColor(255, 244, 206)
WARN_T = RGBColor(120, 68, 10)
WARM_BG = RGBColor(254, 249, 242)

# ── Logo paths (from HubSpot CDN, pre-downloaded) ─────────────────────────
LOGO_PATHS = {
    # ── Original HubSpot logos ────────────────────────────────────────────────
    "LOBO LOGISTICS":                  "/tmp/logos/lobo_logistics.img",
    "MAREX":                           "/tmp/logos/marex.img",
    "BRINKS LAND IMPROVEMENT":         "/tmp/logos/brinks_land.img",
    "EPIC MATERIALS INC":              "/tmp/logos/epic_materials.img",
    "CANTON CONCRETE (DUPLICATE)":     "/tmp/logos/canton_concrete.img",
    "WALKER AG GROUP":                 "/tmp/logos/walker_ag.img",
    "TWIN CITY HAULING":               "/tmp/logos/twin_city.img",
    "STATEWIDE MATERIALS":             "/tmp/logos/statewide.img",
    "PRINCE GEORGE AG":                "/tmp/logos/prince_george.img",
    # ── Enterprise accounts ───────────────────────────────────────────────────
    "AMRIZE: SASK + WINNIPEG":         "/tmp/logos/amrize.img",
    "AMRIZE: NCR-TWIN CITIES":         "/tmp/logos/amrize.img",
    "AMRIZE: GVA (BC)":                "/tmp/logos/amrize.img",
    "AMRIZE: GTA":                     "/tmp/logos/amrize.img",
    "DUFFERIN AGGREGATES (CRH)":       "/tmp/logos/dufferin.img",
    "ZEMBA INC.":                      "/tmp/logos/zemba.img",
    "D CRUPI & SONS, INC.":            "/tmp/logos/d_crupi.img",
    "PJ KEATING CO":                   "/tmp/logos/pj_keating.img",
    # ── Mid-market logos fetched from web ────────────────────────────────────
    "4M TRUCKING":                     "/tmp/logos/4m_trucking.img",
    "BUESING CORP":                    "/tmp/logos/buesing.img",
    "DANIELA TRUCKING & GRADING":      "/tmp/logos/daniela.img",
    "DIAMOND MATERIALS":               "/tmp/logos/diamond.img",
    "FLASH TRUCKING / GOLF AGRONOMICS":"/tmp/logos/flash.img",
    "GERNATT ASPHALT PRODUCTS":        "/tmp/logos/gernatt.img",
    "GEORGE J. IGEL & CO.":            "/tmp/logos/igel.img",
    "GRANITE CONSTRUCTION (SOCAL)":    "/tmp/logos/granite_socal.img",
    "IROQUOIS BAR CORPORATION":        "/tmp/logos/iroquois.img",
    "MANSTEEL REBAR LTD.":             "/tmp/logos/mansteel.img",
    "MMC MATERIALS INC":               "/tmp/logos/mmc_materials.img",
    "N.S. TRUCKING INC.":              "/tmp/logos/ns_trucking.img",
    "PETERSON COMPANIES":              "/tmp/logos/peterson.img",
    "PINERIDGE FARMS INC.":            "/tmp/logos/pineridge.img",
    "QUALITY TRUCKING":                "/tmp/logos/quality_trucking.img",
    "RHINO TRUCKING INC.":             "/tmp/logos/rhino.img",
    "ROCK ON TRUCKS":                  "/tmp/logos/rock_on.img",
    "RPM xCONSTRUCTION":              "/tmp/logos/rpm_x.img",
    "R.W. DUNTEMAN CO.":               "/tmp/logos/rw_dunteman.img",
    "SILVERKING TRUCKING":             "/tmp/logos/silverking.img",
    "TAPANI INC":                      "/tmp/logos/tapani.img",
    "TERRY EQUIPMENT COMPANY":         "/tmp/logos/terry.img",
    "THUNDERBOLT":                     "/tmp/logos/thunderbolt.img",
    # ── Upgraded logos (re-fetched at higher res) ─────────────────────────────
    "TOMLINSON":                       "/tmp/logos/tomlinson.img",
    "TRANS-PHOS INC.":                 "/tmp/logos/trans_phos.img",
    "GULFSHORE TRUCKING LLC":          "/tmp/logos/gulfshore.img",
    "WERDCO BC INC.":                  "/tmp/logos/werdco.img",
    # ── Additional mid-market logos ───────────────────────────────────────────
    "UNITED STATES LIME & MINERALS":   "/tmp/logos/uslm.img",
    "VOLKER STEVIN CONTRACTING":       "/tmp/logos/volker_stevin.img",
    "WESTERN STATES CONTRACTING":      "/tmp/logos/western_states.img",
    "WILLIAMS TRUCKING CO.":           "/tmp/logos/williams.img",
}

# ── Open Intercom support conversations (last 30 days) ────────────────────
INTERCOM_SUPPORT = {
    "AHS": [
        "215474033701302","215474050128786","215474181334923",
        "215473991016651","215473944055565","215473890010511",
    ],
    "AMRIZE: NCR-TWIN CITIES": ["215474022198661","215474197837742"],
    "CEMEX USA":               ["215473856491965","215473866538222"],
    "DANIELA TRUCKING & GRADING": ["215474003597339"],
    "DIAMOND MATERIALS": ["215473992052795","215474228902783"],
    "R.W. DUNTEMAN CO.": ["215474227152056","215474011047073"],
    "IROQUOIS BAR CORPORATION": ["215474231752499"],
    "JW GOLDING": ["215473931769979"],
    "MMC MATERIALS INC": ["215474075885760","215474184671998"],
    "MANSTEEL REBAR LTD.": ["215474224958023"],
    "QUALITY TRUCKING": ["215474197049557","215474213405851"],
    "ROCK ON TRUCKS": [
        "215474239714788","215474228695159","215473996529732","215474131030408",
        "215474229148100","215474197816542","215474215554268","215474215245343",
        "215474151909814","215474142675333","215474083972196","215474003749617",
        "215473850527530","215473853744955",
    ],
    "RONYX LOGISTICS LLC": ["215474205715897"],
    "SILVERKING TRUCKING": ["215473946098787"],
    "STATEWIDE MATERIALS": ["215474230667370","215474216609412"],
    "TAPANI INC": ["215474206679744"],
    "TERRY EQUIPMENT COMPANY": ["215473859191755"],
    "TILCON CT INC": ["215474010729756"],
    "TOMLINSON":     ["215474227880490","215474223429030","215474027908029","215474022268366"],
    "TRANS-PHOS INC.": ["215474227410946","215474196744100","215474007719515"],
}

# ── Granular feature overrides (ranked most→least used/important) ──────────
# Prefixes: Ticket: / Dispatch: / Approvals: / Billing: / Integration: / GPS: / Reports: / Vendor Mgmt:
FEATURE_OVERRIDES = {
    "4M TRUCKING": [
        "Dispatch: central (vendors CANNOT self-dispatch)",
        "Ticket: AI Auto-Scan (active — duplicate detection; error rate concern raised)",
        "Driver App",
        "Fuel Surcharge Add-Ons (URGENT bug TRE-14495 — not cascading to jobs)",
        "Vendor/Contractor Mgmt (King, Gatewood, JNC, RRT, B&C, FGC, WTSC)",
        "Billing → QB CSV export (detailed line items)",
        "GPS Geofence Reports (truck hit count vs. expected ticket count)",
        "Audit Logs (recently enabled)",
        "Approvals (dispatch + accounting teams; unapproval edge cases)",
    ],
    "AHS": [
        "Dispatch: by-load not by-job (avoids GPS shutdown issue)",
        "Driver App (GPS per load)",
        "Permissions / Dispatcher Lite Role (active config)",
        "Billing: Fuel Surcharge built into monthly rate (not a separate add-on)",
        "Ticket Export CSV (leading-zeros bug ongoing)",
        "Reports (Omni) — QB report in progress; vendor name from Site ID pending",
    ],
    "AMRIZE: SASK + WINNIPEG": [
        "Dispatch (to hauler partners — Rock On, M&J, Shaw, Beach Transport)",
        "GPS Tracking (wants round-trip time reports; offline mode needed)",
        "Reports (Omni): wants GPS data added to report builder",
        "Loader App",
    ],
    "BRINKS LAND IMPROVEMENT": [
        "Dispatch (onboarding — was Excel/email; foreman calls truck boss with next-day needs)",
        "Billing: Fuel Surcharge + Standby Time Add-Ons (configuring — ~$0.76/mi; setup trouble)",
        "Billing: multiple rate structures (equipment type + distance + tonnage)",
        "AR/AP Invoices & Settlement",
        "Vendor/Subcontractor Mgmt",
    ],
    "BUESING CORP": [
        "Ticket: Photo Capture + OCR (wants driver reminder system; improved over prior Rocket app)",
        "Driver App",
        "Approvals: two-level payroll review (foreman → second approver before submission)",
        "Timesheets",
    ],
    "CANTON CONCRETE (DUPLICATE)": [
        "Dispatch",
        "Timesheets / Approvals",
        "Integration: QB Enterprise (IIF payroll file — time import via Excel-editable format)",
    ],
    "DANIELA TRUCKING & GRADING": [
        "Dispatch",
        "Driver App",
        "Approvals / Timesheets (Estefania manages — Spanish-speaking team active in Intercom)",
        "Foremen Workflow (wants loader signoff/accountability — not yet built)",
        "Billing → QB Online (QBO) export via Omni (rates not combining per-load — ongoing pain)",
        "Viewer Role (per-site grant — doesn't scale for large orgs)",
    ],
    "GEORGE J. IGEL & CO.": [
        "Ticket: Scanned (majority scanned; some physical remain)",
        "Billing: Fuel Surcharge (manually tracks weekly gov't fuel rates; wants automation)",
        "Vendor/Sub Dispatch (wants sub-of-sub live visibility — who actually shows up on site)",
        "Reports: mobile/field-side reporting, historical + area analysis (replacing Excel)",
    ],
    "GERNATT ASPHALT PRODUCTS": [
        "Dispatch",
        "GPS / Geofence (ticket times + geofence times for billing verification)",
        "Reports (Omni)",
    ],
    "JW GOLDING": [
        "Dispatch",
        "Driver App",
        "Integration: QuickBooks (export format — full integration vs. spreadsheet still TBD)",
    ],
    "MAREX": [
        "Ticket: RFID Auto-Print (tickets auto-print; driver catches through window while in motion)",
        "Ticket: Photo Capture (driver snaps photo → autopopulates ticket data; reduces manual work)",
        "GPS Data (daily report with GPS as proof of payment accuracy)",
        "Integration: QuickBooks Online",
        "Dispatch",
    ],
    "MMC MATERIALS INC": [
        "Dispatch",
        "Ticket: Photo (wants AI auto-scan for ticket# + tonnage — manual entry errors)",
        "Billing: Fuel Surcharge (customers billed higher rate than driver commission — surcharge split)",
        "GPS Tracking",
    ],
    "N.S. TRUCKING INC.": [
        "Dispatch",
        "Driver App",
        "Billing: Fuel Surcharge (manually deducted outside system; added as exact $ add-on in Tread)",
    ],
    "PINERIDGE FARMS INC.": [
        "Approvals (Approve Work + Approve Jobs tabs — active daily use)",
        "Reports: custom driver/dispatch productivity KPIs needed",
        "Integration: QuickBooks (Desktop + Online sync both needed — stalled, no timeline)",
        "Dispatch",
    ],
    "PRIME AGGREGATE TRANSPORTATION": [
        "Dispatch",
        "Billing: Fuel Surcharge (collects surcharges, keeps % margin — bills customer > pays driver)",
        "Driver App",
        "Billing & Settlement",
    ],
    "PRINCE GEORGE AG": [
        "Dispatch",
        "Integration: QuickBooks",
        "Reports: custom daily vendor activity + order status (similar to billing reports)",
        "Driver App",
    ],
    "QUALITY TRUCKING": [
        "Vendor Dispatch: sends 1 dispatch to contractor who manages their own fleet (RPMx is source of 30%+ loads)",
        "Billing: Fuel Surcharge (separated line item on invoices — required for proof/clarity)",
        "Driver App",
        "Reports",
    ],
    "ROCK ON TRUCKS": [
        "Dispatch (central + wants subcontractors to access for owner-operator dispatch)",
        "Billing: Fuel Surcharge Add-Ons (hourly billing; wants weekly-update table with $/% options)",
        "Integration: QB Export (BROKEN — fuel surcharge line items missing; weeks behind on invoicing)",
        "Vendor/Sub Mgmt (dispatches to Amoris Materials + others)",
        "Driver App",
    ],
    "RONYX LOGISTICS LLC": [
        "GPS / Real-Time Truck Location Dashboard (phone-based)",
        "Integration: QuickBooks (report format modified for QB bookkeeping import)",
        "Dispatch",
        "Driver App",
    ],
    "RPM xCONSTRUCTION": [
        "Ticket: AI Auto-Capture (ticket time + pickup time + quantity from photo — no manual entry)",
        "GPS + Geofence Reports (exception reports for misses; truck revenue reports; custom geofence site)",
        "Vendor Dispatch: dispatches to Quality Trucking + Nickel Rock (~300 trucks)",
        "Reports: geofence compliance, daily revenue, rate alignment",
        "Driver App",
    ],
    "R.W. DUNTEMAN CO.": [
        "Ticket: Photo Capture (dirt tickets at dump site — consolidates all records in one place)",
        "Billing: Fuel Surcharge (variable; applied at invoicing not ticket creation; broker surcharges new)",
        "Integration: Vista (payroll import — NOT QuickBooks)",
        "Broker Add-On Charges (fuel surcharge + tolls consolidated at invoice)",
        "Dispatch",
    ],
    "STATEWIDE MATERIALS": [
        "Billing: Fuel Surcharge Automation (8 different surcharges by customer; weekly/monthly update)",
        "Dispatch",
        "Reports (Omni)",
        "Driver App",
    ],
    "TAPANI INC": [
        "Approvals (needs exact job start/end times visible — multi-screen review is pain point)",
        "Telematics / Mileage Reports (state-by-state mileage per vehicle + on-site geofence data needed)",
        "Dispatch",
        "Driver App",
    ],
    "THUNDERBOLT": [
        "Integration: QuickBooks (payroll + customer invoicing — seamless export required)",
        "Timesheets / Driver Time Approvals",
        "Billing: Settlement",
        "Dispatch",
    ],
    "TILCON CT INC": [
        "Approvals: Foreman Signoffs (timestamps + notes — read-only, foreman cannot edit)",
        "Dispatch",
        "Driver App",
        "Reports",
    ],
    "TOP TIER TRUCKING": [
        "Dispatch",
        "Driver App",
        "Integration: QuickBooks (CSV import via QB template — key columns mapped manually)",
    ],
    "TRIO AGGREGATE HAULERS": [
        "Dispatch",
        "Billing: Fuel Surcharge (flat rate; manually calculated — not automated in Tread)",
        "Integration: QuickBooks (QB for billing data import)",
        "Driver App",
    ],
    "TWIN CITY HAULING": [
        "Ticket: Auto-Capture from Photo (wants job#, ticket#, start time, location auto-extracted)",
        "Dispatch",
        "Driver App",
    ],
    "UNITED STATES LIME & MINERALS": [
        "Vendor Mgmt: 9 carriers (Cecil's Transport, Afritz, Tex Sand, BBG, Transwood, JRT, US Transport, Trimac, Thurman)",
        "Dispatch (producer model — manages all hauler vendors centrally)",
        "Driver App (vendor drivers)",
        "Billing & Settlement",
    ],
    "VOLKER STEVIN CONTRACTING": [
        "Approvals: Multi-level (foreman mobile → superintendent office; on-site end-time approval required)",
        "Ticket: Photo Capture (end-of-shift receiving + ticket capture step)",
        "Vendor Onboarding (50-60 vendors in progress — Burnco + others urgently needed)",
        "Dispatch",
        "Timesheets (hourly trucks)",
        "Foreman App (mobile approvals)",
    ],
    "WALKER AG GROUP": [
        "Dispatch + Project Mgmt (active: Sullivan / Plateau Excavation; Ronyx trucking; Lincoln Park Stone materials)",
        "Billing: Fuel Surcharge (market-fluctuation based — dynamic integration TBD; currently unresolved gap)",
        "Vendor Mgmt (Ronyx, US Aggregates, Lincoln Park Stone + others)",
        "Driver App",
    ],
    "WILLIAMS TRUCKING CO.": [
        "Approvals (two-person: verify + approve → AP payables; wants single bulk-approve button)",
        "Dispatch",
        "Driver App",
        "Billing & Settlement",
    ],
}

# ── Usage status & tenure ──────────────────────────────────────────────────
# status: "Primary system" | "Onboarding" | "Sporadic" | "Disengaged" | "Churned"
USAGE_STATUS = {
    "HOLCIM - NORTH CENTRAL (FARGO)":   "Primary system",
    "RAM-CO TRUCKING SERVICES":         "Churned",
    "MIDTEX MATERIALS":                 "Churned",
    "4M TRUCKING":                      "Primary system",
    "AHS":                              "Primary system",
    "AMRIZE: SASK + WINNIPEG":          "Onboarding",
    "ARIZONA AGGREGATE SOLUTIONS":      "Sporadic",
    "BRINKS LAND IMPROVEMENT":          "Primary system",
    "BUESING CORP":                     "Primary system",
    "CANTON CONCRETE (DUPLICATE)":      "Sporadic",
    "CERUTTI & SONS TRANSPORTATION":    "Disengaged",
    "CHARLES H CARTER & SON":           "Primary system",
    "DANIELA TRUCKING & GRADING":       "Primary system",
    "DIAMOND MATERIALS":                "Sporadic",
    "EPIC MATERIALS INC":               "Primary system",
    "FLASH TRUCKING / GOLF AGRONOMICS": "Sporadic",
    "GEORGE J. IGEL & CO.":             "Onboarding",
    "GERNATT ASPHALT PRODUCTS":         "Primary system",
    "GRANITE CONSTRUCTION (SOCAL)":     "Sporadic",
    "IROQUOIS BAR CORPORATION":         "Primary system",
    "JW GOLDING":                       "Primary system",
    "LOBO LOGISTICS":                   "Disengaged",
    "MANSTEEL REBAR LTD.":              "Primary system",
    "MARCC TRUCKING":                   "Onboarding",
    "MAREX":                            "Sporadic",
    "MMC MATERIALS INC":                "Primary system",
    "N.S. TRUCKING INC.":               "Primary system",
    "PETERSON COMPANIES":               "Sporadic",
    "PINERIDGE FARMS INC.":             "Primary system",
    "PRIME AGGREGATE TRANSPORTATION":   "Primary system",
    "PRINCE GEORGE AG":                 "Primary system",
    "QUALITY TRUCKING":                 "Primary system",
    "RHINO TRUCKING INC.":              "Disengaged",
    "ROCK ON TRUCKS":                   "Primary system",
    "RONYX LOGISTICS LLC":              "Primary system",
    "RPM xCONSTRUCTION":               "Primary system",
    "R.W. DUNTEMAN CO.":                "Onboarding",
    "SILVERKING TRUCKING":              "Sporadic",
    "STATEWIDE MATERIALS":              "Primary system",
    "TAPANI INC":                       "Primary system",
    "TERRY EQUIPMENT COMPANY":          "Primary system",
    "THUNDERBOLT":                      "Onboarding",
    "TILCON CT INC":                    "Primary system",
    "TOP TIER TRUCKING":                "Sporadic",
    "TRIO AGGREGATE HAULERS":           "Primary system",
    "TWIN CITY HAULING":                "Disengaged",
    "UNITED STATES LIME & MINERALS":    "Primary system",
    "VOLKER STEVIN CONTRACTING":        "Onboarding",
    "WALKER AG GROUP":                  "Onboarding",
    "WESTERN STATES CONTRACTING":       "Primary system",
    "WILLIAMS TRUCKING CO.":            "Disengaged",
    # Enterprise
    "AMRIZE: NCR-TWIN CITIES":          "Onboarding",
    "AMRIZE: GVA (BC)":                 "Onboarding",
    "AMRIZE: GTA":                      "Primary system",
    "CEMEX USA":                        "Primary system",
    "DUFFERIN AGGREGATES (CRH)":        "Primary system",
    "NATIONAL LIME AND STONE":          "Primary system",
    "TOMLINSON":                        "Primary system",
    "TRANS-PHOS INC.":                  "Primary system",
    "WHITAKER TRANSPORTATION":          "Primary system",
    "ZEMBA INC.":                       "Disengaged",
    # Mid-market new
    "D CRUPI & SONS, INC.":             "Primary system",
    "GULFSHORE TRUCKING LLC":           "Primary system",
    "PJ KEATING CO":                    "Onboarding",
    "R&R TRUCKING, INC.":               "Primary system",
    "UPPAL TRUCKING LTD":               "Disengaged",
    "WERDCO BC INC.":                   "Primary system",
}

# Approximate tenure — derived from contract dates / personality notes
TENURE_APPROX = {
    "4M TRUCKING":                      "~1 yr",
    "AHS":                              "~2 yrs",
    "AMRIZE: SASK + WINNIPEG":          "Pilot",
    "ARIZONA AGGREGATE SOLUTIONS":      "~1 yr",
    "BRINKS LAND IMPROVEMENT":          "~2 yrs",
    "BUESING CORP":                     "~2 yrs",
    "CANTON CONCRETE (DUPLICATE)":      "~1 yr",
    "CERUTTI & SONS TRANSPORTATION":    "~1 yr",
    "CHARLES H CARTER & SON":           "3+ yrs",
    "DANIELA TRUCKING & GRADING":       "~2 yrs",
    "DIAMOND MATERIALS":                "~1 yr",
    "EPIC MATERIALS INC":               "~1 yr",
    "FLASH TRUCKING / GOLF AGRONOMICS": "~2 yrs",
    "GEORGE J. IGEL & CO.":             "< 1 mo",
    "GERNATT ASPHALT PRODUCTS":         "~2 yrs",
    "GRANITE CONSTRUCTION (SOCAL)":     "~1 yr",
    "IROQUOIS BAR CORPORATION":         "~2 yrs",
    "JW GOLDING":                       "~1 yr",
    "LOBO LOGISTICS":                   "~1 yr",
    "MANSTEEL REBAR LTD.":              "~2 yrs",
    "MARCC TRUCKING":                   "< 3 mo",
    "MAREX":                            "~1 yr",
    "MMC MATERIALS INC":                "~2 yrs",
    "N.S. TRUCKING INC.":               "3+ yrs",
    "PETERSON COMPANIES":               "~2 yrs",
    "PINERIDGE FARMS INC.":             "~1 yr",
    "PRIME AGGREGATE TRANSPORTATION":   "~1 yr",
    "PRINCE GEORGE AG":                 "~1 yr",
    "QUALITY TRUCKING":                 "~2 yrs",
    "RHINO TRUCKING INC.":              "~1 yr",
    "ROCK ON TRUCKS":                   "~1 yr",
    "RONYX LOGISTICS LLC":              "~1 yr",
    "RPM xCONSTRUCTION":               "~2 yrs",
    "R.W. DUNTEMAN CO.":                "~1 yr",
    "SILVERKING TRUCKING":              "~1 yr",
    "STATEWIDE MATERIALS":              "~2 yrs",
    "TAPANI INC":                       "~2 yrs",
    "TERRY EQUIPMENT COMPANY":          "~1 yr",
    "THUNDERBOLT":                      "< 3 mo",
    "TILCON CT INC":                    "~2 yrs",
    "TOP TIER TRUCKING":                "~2 yrs",
    "TRIO AGGREGATE HAULERS":           "~1 yr",
    "TWIN CITY HAULING":                "~1 yr",
    "UNITED STATES LIME & MINERALS":    "~2 yrs",
    "VOLKER STEVIN CONTRACTING":        "< 3 mo",
    "WALKER AG GROUP":                  "< 1 mo",
    "WESTERN STATES CONTRACTING":       "~2 yrs",
    "WILLIAMS TRUCKING CO.":            "~2 yrs",
    # Enterprise
    "AMRIZE: NCR-TWIN CITIES":          "~2 yrs",
    "AMRIZE: GVA (BC)":                 "3+ yrs",
    "AMRIZE: GTA":                      "3+ yrs",
    "CEMEX USA":                        "3+ yrs",
    "DUFFERIN AGGREGATES (CRH)":        "3+ yrs",
    "NATIONAL LIME AND STONE":          "3+ yrs",
    "TOMLINSON":                        "3+ yrs",
    "TRANS-PHOS INC.":                  "3+ yrs",
    "WHITAKER TRANSPORTATION":          "3+ yrs",
    "ZEMBA INC.":                       "~2 yrs",
    # Mid-market new
    "D CRUPI & SONS, INC.":             "3+ yrs",
    "GULFSHORE TRUCKING LLC":           "3+ yrs",
    "PJ KEATING CO":                    "3+ yrs",
    "R&R TRUCKING, INC.":               "~2 yrs",
    "UPPAL TRUCKING LTD":               "3+ yrs",
    "WERDCO BC INC.":                   "3+ yrs",
    "HOLCIM - NORTH CENTRAL (FARGO)":  "~1 yr",
}

# ── Contract renewal dates (from Lovable CS Tracker) ──────────────────────
RENEWAL_DATES = {
    "MAREX":                           "2026-06-29",
    "GRANITE CONSTRUCTION (SOCAL)":    "2026-08-30",
    "TRANS-PHOS INC.":                 "2026-09-13",
    "EPIC MATERIALS INC":              "2026-11-22",
    "PRINCE GEORGE AG":                "2026-11-30",
    "GERNATT ASPHALT PRODUCTS":        "2026-12-30",
    "WERDCO BC INC.":                  "2027-01-31",
    "WHITAKER TRANSPORTATION":         "2027-01-31",
    "HOLCIM - NORTH CENTRAL (FARGO)":  "2027-03-02",
    "PETERSON COMPANIES":              "2027-03-08",
    "SILVERKING TRUCKING":             "2027-03-08",
    "WESTERN STATES CONTRACTING":      "2027-03-11",
    "D CRUPI & SONS, INC.":            "2027-03-30",
    "TAPANI INC":                      "2027-04-02",
    "CANTON CONCRETE (DUPLICATE)":     "2027-04-20",
    "MANSTEEL REBAR LTD.":             "2027-05-30",
    "RPM xCONSTRUCTION":               "2027-07-30",
    "TOMLINSON":                       "2027-08-24",
    "QUALITY TRUCKING":                "2027-09-20",
    "IROQUOIS BAR CORPORATION":        "2027-09-29",
    "PRIME AGGREGATE TRANSPORTATION":  "2028-02-28",
}

# ── Company data ───────────────────────────────────────────────────────────
midmarket_companies = [
  {
    "name": "4M TRUCKING",
    "health": "yellow",   # ← updated from green: COO frustrated, URGENT bug open, neutral impression May 7
    "customer_type": "Hauler",
    "arr": "$52,800",
    "trucks": "75",
    "location": "TX",
    "csm": "Latefa Redjouh",
    "owner": "Adam Murray",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/39582871444",
    "what": "75-truck hauler. Dispatch controlled by 4M; sub-vendor contractors upload tickets. Separate dispatch team and invoicing team.",
    "connects_with": "Sub-vendors: King, Gatewood, JNC, RRT, B&C, FGC, WTSC (brokerage-type contractors they dispatch to)",
    "main_contacts": [("Lonie Packer", "IT/Ops — day-to-day Tread admin"), ("Wheeler Renshaw", "CFO"), ("Linda Stamps", "COO — key stakeholder, frustrated, absent from calls")],
    "tread_features": ["Dispatch (central control — vendors cannot self-dispatch)", "Driver App", "Vendor / Contractor Management", "Settlement & Invoicing", "Fuel Surcharge Add-Ons", "Audit Logs (recently enabled)"],
    "personality": "May 7 call: overall impression neutral. COO Linda has 'had enough' with data inconsistencies and was absent again. Lonie is cautiously optimistic but hasn't ruled out alternatives. Email confirmation feature (in dev) was very well received — could help rebuild trust.",
    "activity": [
      ("May 7", "HeySam", "Check-in: fuel surcharge add-on bug open (not propagating to jobs); audit logs enabled; email receipt feature demoed — Lonie loved it; COO absent again"),
      ("May 7", "Intercom", "Driver app notification question (open); phone conflict (resolved)"),
    ],
    "tickets": [
      ("TRE-14495", "URGENT — Order Add-Ons not cascading to jobs (fuel surcharge bug)"),
      ("OPS-610", "PROJ-55800 missing materials — Lonie reported project data changed unexpectedly"),
    ],
    "risks": ["COO Linda Stamps frustrated & absent — account-level trust at risk. URGENT ticket TRE-14495 open. Data integrity unexplained changes are core pain. No ARR recorded despite 75 trucks."],
    "systems": [],
  },
  {
    "name": "AHS",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$39,000",
    "trucks": "100",
    "location": "Ocala, FL",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/33696058900",
    "what": "Aggregate hauling company, 100 trucks. Active in permissions config and reporting.",
    "connects_with": "Producers and quarries (material sources); construction sites (delivery destinations)",
    "main_contacts": [("Katie Stewart", "Admin Assistant"), ("Ken McDonald", "Ops"), ("Johnathan Rosa", "")],
    "tread_features": ["Dispatch", "Reporting (Omni)", "Permissions / Dispatcher Lite Role", "Ticket Export (CSV)"],
    "personality": "Technically engaged, active in QBRs. CSV leading-zeros bug is a recurring frustration. Responds well to structured sessions.",
    "activity": [
      ("Apr 30", "HeySam", "Permissions session — Dispatcher Lite scoped; ticket export confirmed; leading-zeros CSV bug"),
      ("May 8", "Gmail", "QB Report Update — 3 of 4 items resolved; vendor name from Site ID still pending"),
    ],
    "tickets": [("REP-126", "Report: scheduled vs. delivered loads & tons daily"), ("MBL-1156", "Mobile app error for driver 118A Maikel Sosa")],
    "risks": [],
    "systems": [],
  },
  {
    "name": "ARIZONA AGGREGATE SOLUTIONS",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$21,420",
    "trucks": "—",
    "location": "Phoenix, AZ",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/53211375157",
    "what": "Aggregate producer and supplier in Phoenix, AZ.",
    "connects_with": "Haulers (specific vendors unknown), construction companies",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Order Management"],
    "personality": "Quiet customer — no recent meetings or conversations found.",
    "activity": [],
    "tickets": [],
    "risks": ["No recent conversations found."],
    "systems": [],
  },
  {
    "name": "BRINKS LAND IMPROVEMENT",
    "health": "green",
    "customer_type": "Construction",
    "arr": "$20,250",
    "trucks": "—",
    "location": "Pleasant Hill, OR",
    "csm": "unassigned",
    "owner": "JP Pasteur (inactive)",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/52742127746",
    "what": "Land improvement and earthwork in OR. Fuel surcharge per-mile not natively supported.",
    "connects_with": "Haulers and material producers (specific vendors not identified yet)",
    "main_contacts": [("John Redden", "Primary contact — wants in-person sessions")],
    "tread_features": ["Dispatch", "Freight Rates + Fuel Surcharge", "Approvals", "Driver CSV Import"],
    "personality": "Relationship-oriented. John Redden strongly prefers in-person sessions over remote. Hands-on and curious about the platform.",
    "activity": [("May 5", "HeySam", "Rates session — freight rates, fuel surcharge, approvals; John wants in-person sessions")],
    "tickets": [("TRE-14142", "Excel corrupts +1 phone prefix on CSV save (422 errors)"), ("TRE-14141", "Address not required in vendor CSV — silent 422 errors")],
    "risks": ["Owner JP Pasteur is inactive — needs reassignment."],
    "systems": [],
  },
  {
    "name": "BUESING CORP",
    "health": "green",
    "customer_type": "Construction",
    "arr": "$59,800",
    "trucks": "—",
    "location": "Phoenix, AZ",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409519445",
    "what": "55+ year leader in excavation, backfill, shoring, shotcrete, trucking, renewables in AZ. Phase 2 (AR/invoicing) underway.",
    "connects_with": "Hauler subcontractors and material vendors (specific names not in Tread data)",
    "main_contacts": [("Bryon Matesi", "President/COO"), ("Ravi Ramachandran", "CIO"), ("Tyler Steers", "Dispatch Manager"), ("Nicki Click", "Payroll")],
    "tread_features": ["Dispatch", "AR / Invoicing (Phase 2)", "Bulk Ticket Download", "Contract Item Mapping", "Reporting"],
    "personality": "Enterprise sophistication. CIO (Ravi) drives tech adoption. Multiple departments engaged — dispatch, payroll, IT. Strong account.",
    "activity": [("Mar 26", "HeySam", "Phase 2 AR kick-off — AR extract format, bulk ticket download, contract item mapping")],
    "tickets": [("REP-191", "Payroll dashboard report — contact Nicki Click")],
    "risks": [],
    "systems": [],
  },
  {
    "name": "CANTON CONCRETE (DUPLICATE)",
    "health": "gray",
    "customer_type": "Construction",
    "arr": "—",
    "trucks": "—",
    "location": "—",
    "csm": "Latefa Redjouh",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/52723975493",
    "what": "Concrete company. Two HubSpot records for same domain — needs merge. Strong driver adoption from training.",
    "connects_with": "Material suppliers, construction sites",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Driver App", "GPS Permissions", "Rain-Day Dispatch"],
    "personality": "Driver adoption clicked in training — GPS and job flows working well. CRM data is messy. No CSM assigned.",
    "activity": [("May 1", "HeySam", "Driver training — GPS permissions, job flows, rain-day dispatch; strong adoption")],
    "tickets": [],
    "risks": ["Duplicate HubSpot records (IDs 52723975493 and 52195729256). No CSM assigned."],
    "systems": [],
  },
  {
    "name": "CERUTTI & SONS TRANSPORTATION",
    "health": "gray",
    "customer_type": "Hauler",
    "arr": "$20,000",
    "trucks": "150",
    "location": "Fresno, CA",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/21311496857",
    "what": "Freight and aggregate hauler, 150 trucks. Last contact Aug 2025 — 8+ months of silence.",
    "connects_with": "Material producers and construction companies (as hauler); used Samsara + Fleetio prior",
    "main_contacts": [("Connor", "Ops contact — last request: remove all equipment from setup")],
    "tread_features": ["Dispatch"],
    "personality": "Silent for 8+ months. May have already churned. No CSM to re-engage.",
    "activity": [],
    "tickets": [("TRE-9491", "Connor requested all equipment removed from setup")],
    "risks": ["No CSM. Last contact Aug 2025 — 8+ months of silence. High churn risk."],
    "systems": ["Custom ERP", "Samsara", "Fleetio"],
  },
  {
    "name": "CHARLES H CARTER & SON",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$36,000",
    "trucks": "—",
    "location": "Fairfield, IL",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409543803",
    "what": "General commodity hauler in Southern IL — rock, sand, gravel, dirt, asphalt, grain, equipment.",
    "connects_with": "Producers, quarries, grain facilities, construction sites in IL/surrounding region",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Load Tracking", "Settlement"],
    "personality": "Long-tenured, stable, low-maintenance. Contract expired Jan 2024 but still active and paying. 30 HubSpot contacts.",
    "activity": [],
    "tickets": [],
    "risks": ["Contract expired Jan 2024 — still active."],
    "systems": [],
  },
  {
    "name": "DANIELA TRUCKING & GRADING",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$10,800",
    "trucks": "12",
    "location": "FL",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/33697686122",
    "what": "Small trucking & grading company in FL, 12 trucks. Operates under Kerrigan Holdings. Uses Tread for dispatch, driver app, foremen workflow, timesheets, and QuickBooks export.",
    "connects_with": "Construction sites and material producers in FL (specific Tread connections unknown — check with Adam Murray)",
    "main_contacts": [("Amy", "Primary dispatcher"), ("Estefania Omana", "Approvals / timesheets — active in Intercom (Spanish-speaking)"), ("Daniela Paunache", "Stakeholder / management escalation contact")],
    "tread_features": ["Dispatch", "Driver App", "Foremen Workflow", "Timesheets / Approvals", "QuickBooks Export (Omni)", "Viewer Role"],
    "personality": "VERY VOCAL. Management level has escalated repeatedly over reporting accuracy — months of broken promises on delivery rate CSV fix. Higher-ups explicitly not satisfied. Adam Murray is CSM. Spanish-speaking staff in Intercom (Estefania). Will email order numbers when issues occur — detail-oriented and persistent.",
    "activity": [
      ("Apr 21", "Intercom", "Estefania (ES) asked how to manually edit missed lunch break in Approvals — tagged Priority + Feature Request, linked to CUS-116, still open"),
      ("Apr 13", "Linear", "REP-223 filed — rates multiplied per job not per load in report for ORD-664460 (resolved May 6)"),
      ("Apr 2026", "HeySam", "Management escalation: delivery rate CSV issue raised with Tim & Adam, told it needs QB integration to fix — higher-ups not satisfied, citing months of prior promises"),
    ],
    "tickets": [
      ("REP-13", "URGENT — Driver Performance Report (In Review since Sep 30 — overdue)"),
      ("PRO-1354", "QuickBooks export: combine Material+Freight into one line per load (Backlog, updated May 7)"),
      ("PRO-1423", "Department not cascading from Project to Order consistently — breaks billing"),
      ("CUS-116", "Allow manual editing of missed lunch breaks in Approvals (Backlog, Apr 21)"),
      ("OPS-557", "Jan 6 feedback: QBO import mapping, driver rate calc errors, freight total columns empty (Triage)"),
    ],
    "risks": ["Management actively frustrated — reporting accuracy issues unresolved for months, prior promises broken. URGENT REP-13 driver report overdue since Sept 30. Multiple backlog items with no timeline."],
    "systems": ["QuickBooks (QBO export via Omni)"],
  },
  {
    "name": "DIAMOND MATERIALS",
    "health": "gray",
    "customer_type": "Construction",
    "arr": "$37,620",
    "trucks": "—",
    "location": "Wilmington, DE",
    "csm": "unassigned",
    "owner": "Adam Murray (inactive)",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9578182030",
    "what": "Full-service site contractor — excavation + site work from backhoes to 35 CY scrapers.",
    "connects_with": "CNT Trucking (vendor they manage in Tread); Tri-County Materials (local partner)",
    "main_contacts": [("Liz Toto", ""), ("Harvey Snow", ""), ("Ed Holston", "Sales Manager"), ("Lori Kibler", "Open Intercom contact")],
    "tread_features": ["Notifications", "Invoicing", "Driver Management"],
    "personality": "Historically self-sufficient but notifications broke Jan–Mar 2026 and went unaddressed for 3 months. Risk of quiet disengagement. No CSM to own the relationship.",
    "activity": [
      ("May 5", "HeySam", "Intro with Latefa — notifications broken Jan-Mar (unreported until Apr 21); invoicing interest"),
      ("Apr 20", "Intercom", "Lori Kibler: notification behavior question — still open"),
      ("May 8", "Gmail", "Driver add request for CNT Trucking vendor account"),
    ],
    "tickets": [],
    "risks": ["No CSM. Owner inactive. Notification bug Apr 21 — untracked in Linear. Intercom open."],
    "systems": [],
  },
  {
    "name": "EPIC MATERIALS INC",
    "health": "green",
    "customer_type": "Producer",
    "arr": "$22,000",
    "trucks": "10",
    "location": "FL",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/40711295642",
    "what": "Aggregate + crushing in FL. Also known as 'Krushin It'. 10 trucks + mobile crushing unit.",
    "connects_with": "Freight carriers: Fraser Hauling, SE Rentals; construction contractors in FL",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Order Management"],
    "personality": "Small, dual-brand (Epic + Krushin It), healthy and quiet. No open issues.",
    "activity": [],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
  {
    "name": "FLASH TRUCKING / GOLF AGRONOMICS",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$15,300",
    "trucks": "—",
    "location": "LaBelle, FL",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/17752433706",
    "what": "Flash Trucking hauls for Golf Agronomics — custom soil blending for top golf courses in FL, GA, NC, SC.",
    "connects_with": "Golf Agronomics (primary customer); golf courses in FL, GA, NC, SC",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Load Tracking"],
    "personality": "Very niche use case. Contract expired Mar 2025, status unclear. Self-managing.",
    "activity": [],
    "tickets": [],
    "risks": ["No CSM. Contract expired Mar 2025. Still listed as Onboarding — status unclear."],
    "systems": [],
  },
  {
    "name": "GEORGE J. IGEL & CO.",
    "health": "gray",
    "customer_type": "Construction",
    "arr": "$130,000",
    "trucks": "—",
    "location": "Columbus, OH",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/52913103758",
    "what": "Major construction company in central OH — earthwork, site dev, concrete, utilities, paving, deep foundations. 90-day pilot starts May 11.",
    "connects_with": "Haulers, material suppliers, subcontractors in OH (vendor setup in progress)",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Sage 300 CSV Integration (in setup — $15K)"],
    "personality": "Brand-new pilot. Understaffed — foreman resistance is the key risk. Needs white-glove onboarding.",
    "activity": [("Apr 27", "HeySam", "Contract finalized — 90-day pilot May 11; Sage 300 CSV integration; $15K fee")],
    "tickets": [],
    "risks": ["CRITICAL ONBOARDING: Pilot starts May 11. Understaffed — foreman resistance is key risk."],
    "systems": ["Sage 300"],
  },
  {
    "name": "GERNATT ASPHALT PRODUCTS",
    "health": "green",
    "customer_type": "Producer",
    "arr": "$60,000",
    "trucks": "—",
    "location": "Collins, NY",
    "csm": "Latefa Redjouh",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9282534708",
    "what": "Part of Gernatt Family of Companies — 50+ years in sand, gravel, asphalt in Western NY. Wind farm hauler onboarding underway.",
    "connects_with": "Hired haulers (third-party vendors); wind farm operators; construction contractors",
    "main_contacts": [("Amanda Fisher", "Salesperson"), ("Jason Hopkins", ""), ("R. Lock", ""), ("Salvatore", "")],
    "tread_features": ["Dispatch", "FileMaker Integration", "Reporting (Omni)", "Route Drawing (interest)", "Wind Farm Hauler Onboarding"],
    "personality": "50+ year family company with enterprise sophistication. FileMaker integration refinement ongoing. Route drawing is a wishlist feature. Engaged and collaborative.",
    "activity": [("Apr 30", "HeySam", "Account update — wind farm hauler onboarding, FileMaker refinement, route drawing interest")],
    "tickets": [("TRE-8689", "Slow scrolling + report spinning"), ("OPS-434", "Driver usage / compliance dashboard not loading"), ("TRE-9412", "500 error after duplicating order ORD-208468")],
    "risks": [],
    "systems": ["FileMaker"],
  },
  {
    "name": "IROQUOIS BAR CORPORATION",
    "health": "green",
    "customer_type": "Construction",
    "arr": "$35,207",
    "trucks": "—",
    "location": "Lackawanna, NY",
    "csm": "Latefa Redjouh",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409544027",
    "what": "Native American-owned full-service construction in Western NY — concrete, rebar, trucking, demolition across four divisions.",
    "connects_with": "Haulers and material suppliers (multi-division, specific names not identified)",
    "main_contacts": [("Mike Hatsios", ""), ("Jenna Zawodzinski", "Escalated urgent dispatch issue — 37+ days open")],
    "tread_features": ["Dispatch", "Geofencing", "Permissions / Roles", "Order Management"],
    "personality": "Multi-division operation with complex workflows. Active on platform but frustrated by multiple open bugs. Jenna escalated an urgent issue 37+ days ago.",
    "activity": [
      ("Apr 8", "HeySam", "Sync — permissions confusion, duplicate site bug, pagination bug, material duplication; geofence ghost-dispatch fixed"),
      ("Apr 1", "Intercom", "Urgent dispatch issue from Jenna Zawodzinski — OPEN 37+ days"),
    ],
    "tickets": [],
    "risks": ["IMMEDIATE RISK: Dispatch issue open 37+ days (Intercom). Enroute bug open since Aug 2025."],
    "systems": [],
  },
  {
    "name": "JW GOLDING",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$56,700",
    "trucks": "—",
    "location": "—",
    "csm": "unassigned",
    "owner": "Adam Murray",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/43858562940",
    "what": "Construction/trucking company. Went live Dec 2025. Driver payout was THE reason they chose Tread.",
    "connects_with": "Material producers and construction sites (specific vendors not identified)",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Rate Management", "Driver Payout (promised — delayed to Sept 30)"],
    "personality": "Trust eroding rapidly. Driver payout promised 3 months ago, now delayed to Sept 30. App sent 15 trucks to a driver's living room. Needs urgent wins.",
    "activity": [("Apr 21", "HeySam", "Dispatch review — rates not cascading, $10K invoice shortfalls, app sent 15 trucks to driver's home")],
    "tickets": [("PRO-1628", "Drivers confusing attachment CTAs with scale ticket upload")],
    "risks": ["HIGH RISK: Driver payout (promised 3 months ago) delayed to Sept 30. Trust eroding fast."],
    "systems": [],
  },
  {
    "name": "LOBO LOGISTICS",
    "health": "red",
    "customer_type": "Hauler",
    "arr": "$75,600",
    "trucks": "—",
    "location": "Aurora, CO",
    "csm": "unassigned",
    "owner": "unassigned",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/35139437235",
    "what": "Logistics company in Aurora, CO. $75.6K ARR — high-value account with zero recent engagement.",
    "connects_with": "Producers and construction companies in CO (specific connections unknown)",
    "main_contacts": [],
    "tread_features": ["Unknown — no recent contact to verify"],
    "personality": "Complete silence. No meetings, emails, or Intercom. $75.6K ARR at serious churn risk. Needs immediate outreach.",
    "activity": [],
    "tickets": [],
    "risks": ["RED HEALTH. No HubSpot owner. Zero meetings or emails found — silent churn risk at $75.6K ARR."],
    "systems": [],
  },
  {
    "name": "MANSTEEL REBAR LTD.",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$20,800",
    "trucks": "10",
    "location": "Richmond Hill, ON",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/30636557060",
    "what": "One of Canada's largest rebar suppliers. Major infrastructure + residential projects across Ontario. ELD (Azuga) integration active.",
    "connects_with": "Construction contractors in Ontario (as rebar supplier + delivery); LP ERP system",
    "main_contacts": [],
    "tread_features": ["Dispatch", "RESTful API Integration (LP ERP)", "Driver App", "ELD / Azuga Integration"],
    "personality": "Technology-forward — automated 2 PM daily delivery transfers via API. ELD connectivity is an active pain point. Sameer personally involved = escalation signal.",
    "activity": [
      ("Mar 30", "HeySam", "LP integration planning — automated 2 PM daily delivery transfers via RESTful API"),
      ("May 8", "Gmail", "Drivers can't connect to trucks via tablets/phones — Sameer escalated to Azuga (ELD)"),
      ("May 7", "Intercom", "Contract + support thread open; Sameer looped in"),
    ],
    "tickets": [],
    "risks": ["No CSM. Driver ELD connectivity issue active (Sameer/CEO involved)."],
    "systems": ["LP (on-premise ERP)"],
  },
  {
    "name": "MARCC TRUCKING",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$36,000",
    "trucks": "70",
    "location": "Dallas, TX",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/50593255295",
    "what": "70-truck hauler in Dallas, TX. Recently onboarded — focused on geofence setup and vendor accounts.",
    "connects_with": "Producers and construction companies in TX (vendor accounts being set up)",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Polygon Geofencing", "Vendor Management", "Driver App", "Phone Deduplication"],
    "personality": "Recently onboarded, actively configuring. Clean engagement — no open issues.",
    "activity": [("Apr 22", "HeySam", "Onboarding — polygon geofences, vendor/driver accounts, duplicate phone cleanup")],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
  {
    "name": "MAREX",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$50,000",
    "trucks": "450",
    "location": "—",
    "csm": "Latefa Redjouh",
    "owner": "Jonathan Luke (inactive)",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/33372126639",
    "what": "Freight brokerage with 450 trucks — largest fleet in portfolio. Still evaluating Tread vs. Turvo.",
    "connects_with": "Construction companies and shippers (as freight broker); previously using Turvo",
    "main_contacts": [("Hayden", "Key decision-maker — currently absent")],
    "tread_features": ["Dispatch", "PDF Rate Management", "QuickBooks Integration (target)", "Approvals"],
    "personality": "Not yet committed to Tread — actively evaluating Turvo as alternative. Key decision-maker (Hayden) absent. Inactive owner needs reassignment. Highest-risk large fleet.",
    "activity": [("Mar 31", "HeySam", "Account update — NOT committed to Tread vs. Turvo; PDF/rate feature live; QB integration delayed")],
    "tickets": [("TRE-9575", "Ended jobs not showing in Approvals tab")],
    "risks": ["AT RISK: Still evaluating Turvo. Key decision-maker absent. Inactive owner."],
    "systems": ["Turvo (evaluating)"],
  },
  {
    "name": "MMC MATERIALS INC",
    "health": "green",
    "customer_type": "Producer",
    "arr": "$124,687",
    "trucks": "100",
    "location": "Ridgeland, MS",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409226294",
    "what": "Leading ready-mix concrete producer (MS, TN, AR, LA). Est. 1927. Also known as MS Aggregate Haulers. 100 trucks.",
    "connects_with": "Hauler vendors (dispatched externally); construction contractors across 4 states",
    "main_contacts": [("Michael Weldon", "GM of Logistics"), ("Scott Craft", "Division President"), ("Bobby Dowdy", "VP Technical Services"), ("Christy Jones", "")],
    "tread_features": ["Dispatch", "Mobile Driver App", "Account Flattening (in progress)", "Shift Hours Reporting", "GPS Tracking"],
    "personality": "Complex multi-region account. Division President (Scott Craft) engaged. Account flattening = drivers can't see rates on mobile until complete. Data-focused team.",
    "activity": [
      ("Apr 27", "HeySam", "Account flattening proposal — drivers can't see rates until complete; shift hours report in dev"),
      ("Apr 19", "Intercom", "Driver GPS/connectivity issue — open"),
    ],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
  {
    "name": "N.S. TRUCKING INC.",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$36,000",
    "trucks": "—",
    "location": "Cocoa, FL",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409418529",
    "what": "~20 years of aggregate hauling in Central FL and the Space Coast.",
    "connects_with": "Aggregate producers and quarries (material sources); construction sites on the Space Coast",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Settlement", "Vendor Management", "Truck Number Reconciliation"],
    "personality": "Long-tenured, experienced hauler. Settlement accuracy is the core concern — truck number discrepancies between vendor vs. driver views. Latefa recently introduced as new CSM.",
    "activity": [("May 6", "HeySam", "Settlement troubleshooting — truck number discrepancy vendor vs. driver; Latefa introduced as new CSM")],
    "tickets": [],
    "risks": ["No ARR recorded."],
    "systems": [],
  },
  {
    "name": "PETERSON COMPANIES",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$10,000",
    "trucks": "—",
    "location": "Chisago City, MN",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/23903933479",
    "what": "Outdoor contractor in Upper Midwest — excavation, irrigation, landscaping, athletic field installation.",
    "connects_with": "M&J Trucking (partner — building their own software); other local haulers and suppliers",
    "main_contacts": [],
    "tread_features": ["Dispatch (limited / testing)"],
    "personality": "Paying for software they're not fully using. Broker pool shrinking. M&J Trucking (key partner) building own software. Ownership transition in progress.",
    "activity": [("Feb 27", "HeySam", "Resuming Tread testing; M&J building own software; ownership transition underway")],
    "tickets": [],
    "risks": ["Paying for unused software. Broker pool shrinking. M&J building own software."],
    "systems": [],
  },
  {
    "name": "PINERIDGE FARMS INC.",
    "health": "yellow",
    "customer_type": "Mixed",
    "arr": "$48,600",
    "trucks": "100",
    "location": "Honolulu, HI",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/20182904152",
    "what": "Hauling + aggregate in HI since 1987. Construction trucking, aggregates, mobile crushing and screening.",
    "connects_with": "Construction companies and contractors in Hawaii (island market)",
    "main_contacts": [],
    "tread_features": ["Dispatch", "QuickBooks Sync (pending — top priority)", "Mobile Crushing/Screening", "Settlement"],
    "personality": "Island-market account — unique logistics constraints. QuickBooks sync is #1 priority with no timeline. Support continuity concern after previous CSM departure.",
    "activity": [("Apr 16", "HeySam", "QB sync still pending (top priority); support continuity concern; contract renewal imminent")],
    "tickets": [],
    "risks": ["QB sync has no timeline. Support continuity concern."],
    "systems": ["Verizon Works", "QuickBooks", "Excel"],
  },
  {
    "name": "PRIME AGGREGATE TRANSPORTATION",
    "health": "yellow",
    "customer_type": "Mixed",
    "arr": "$83,160",
    "trucks": "—",
    "location": "TX",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/23904946598",
    "what": "TX aggregate supplier + construction. Owns quarries (crushed stone, sand, gravel). Site prep, excavation, grading. Hauls own materials.",
    "connects_with": "Construction companies and contractors in TX (as producer + self-haul operation)",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Settlement", "Fuel Surcharge Tables", "Motive Integration"],
    "personality": "Complex: owns quarries AND hauls. Fuel surcharge managed manually (pending AI automation). Settlement training recently completed.",
    "activity": [("May 1", "HeySam", "Fuel surcharge + settlement training; fuel price index manual until AI automation built")],
    "tickets": [],
    "risks": [],
    "systems": ["Motive"],
  },
  {
    "name": "PRINCE GEORGE AG",
    "health": "green",
    "customer_type": "Agriculture",
    "arr": "$63,000",
    "trucks": "20",
    "location": "—",
    "csm": "Latefa Redjouh",
    "owner": "unassigned",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/35427902307",
    "what": "Agricultural trucking company, 20 trucks.",
    "connects_with": "Agricultural producers, farmers, grain facilities",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Driver App"],
    "personality": "Quiet, healthy account. No HubSpot owner. Last meeting Jan 22 — outside 90-day window. Needs a check-in.",
    "activity": [],
    "tickets": [],
    "risks": ["No HubSpot owner. Last meeting Jan 22 — outside 90-day window."],
    "systems": [],
  },
  {
    "name": "QUALITY TRUCKING",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$150,000",
    "trucks": "—",
    "location": "Little Rock, AR",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/35274092586",
    "what": "Large aggregate hauler in AR. Sub-hauler tightly integrated with RPM xConstruction. 3rd-highest ARR.",
    "connects_with": "RPM xConstruction (sends 30%+ of their jobs); receives orders directly from RPMx in Tread",
    "main_contacts": [("Charles Schmidt", ""), ("Britney Richards", ""), ("Slade Zeigler", ""), ("D'Anne Temple", "")],
    "tread_features": ["Dispatch", "Reporting (Omni)", "Split Tickets", "Settlement", "Driver App", "RPMx Order Integration"],
    "personality": "Highly integrated with RPMx — Quality Trucking follows RPMx's platform choices. Technically active, reporting-focused. Key insight: if RPMx churns, Quality Trucking follows.",
    "activity": [("May 6", "Intercom", "Driver phone number conflict — resolved")],
    "tickets": [
      ("TRE-14087", "Split action hidden on RPMx orders with vendors_can_split=true, qty=1"),
      ("MBL-1797", "Driver app: no connection off WiFi (works in Safari)"),
      ("REP-115", "Material, vendor, customer costs report for internal drivers"),
    ],
    "risks": [],
    "systems": [],
  },
  {
    "name": "RHINO TRUCKING INC.",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$1,200",
    "trucks": "—",
    "location": "Indianapolis, IN",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409450550",
    "what": "Transportation company in Indianapolis. Lowest ARR in portfolio. Contract expired Dec 2024.",
    "connects_with": "Producers and construction sites in IN (specific names unknown)",
    "main_contacts": [],
    "tread_features": ["Dispatch"],
    "personality": "Likely near-churned. No CSM, expired contract, last contact Nov 2025. Low ARR makes re-engagement economics difficult.",
    "activity": [],
    "tickets": [],
    "risks": ["No CSM. Contract expired Dec 2024. $2,040 ARR. Last contact Nov 2025."],
    "systems": [],
  },
  {
    "name": "ROCK ON TRUCKS",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$29,160",
    "trucks": "100",
    "location": "Waite Park, MN",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409480104",
    "what": "Premier aggregate hauler in MN. Dedicated, OTR, short + long haul. 100 trucks.",
    "connects_with": "Material source: Amoris Materials (they haul from Amoris yard, then deliver to construction sites)",
    "main_contacts": [("Krystal Vierkant", "CEO — personally involved and skeptical")],
    "tread_features": ["Dispatch", "Invoicing (PDF)", "Reporting", "Split Hours", "Settlement"],
    "personality": "CEO Krystal directly engaged: 'quite nervous and not trusting the system.' Sameer (CEO of Tread) personally reviewing. Multiple active Intercom + Gmail issues. Trust-rebuilding is urgent.",
    "activity": [
      ("May 9", "Gmail", "Krystal (CEO): 'We are quite nervous and honestly not trusting the system.' Sameer personally reviewing."),
      ("May 2026", "Intercom", "4 open: invoice PDF error, vendor name fix, canceled jobs display, split hours error"),
    ],
    "tickets": [("REP-274", "Report: dates, ticket numbers, truck numbers, quantities billed per project")],
    "risks": ["IMMEDIATE RISK: CEO trust breakdown. Sameer personally involved. 4 Intercom + 3 Gmail active issues."],
    "systems": [],
  },
  {
    "name": "RONYX LOGISTICS LLC",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$58,800",
    "trucks": "50",
    "location": "Multi-state (incl. Chicago, IL)",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/49562040681",
    "what": "Rapidly scaling logistics company — 200 trucks in one state, +100 more in May, 5 states including Chicago.",
    "connects_with": "Construction companies and producers across 5 states (as hauler/broker at scale)",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Vendor Management", "Document / Compliance Management", "RBAC (supervisor rate hide)", "Scaling Ops"],
    "personality": "High-growth, tech-savvy. Scaling fast — RBAC and supervisor permissions are current friction. 206-reply vendor management thread shows operational complexity.",
    "activity": [
      ("Apr 15", "HeySam", "Compliance docs discussion — Document Types vs. RIMS; RBAC gap for supervisor rate visibility"),
      ("May 5", "Intercom", "Driver assignment issue ORD-826804 (open); vendor management thread (206 replies)"),
    ],
    "tickets": [("TRE-14090", "Ability to hide rate from supervisor/foreman roles")],
    "risks": [],
    "systems": [],
  },
  {
    "name": "R.W. DUNTEMAN CO.",
    "health": "gray",
    "customer_type": "Construction",
    "arr": "$25,916",
    "trucks": "75",
    "location": "Addison, IL",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/10627140958",
    "what": "Construction company in IL, 75 trucks. Vista ERP integration actively being built.",
    "connects_with": "Haulers, material vendors, subcontractors in IL",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Vista ERP Export Integration", "Approvals", "Driver Start-Time Notifications"],
    "personality": "ERP-integration focused — Vista export session deep in progress. Two open Intercom threads show active engagement. Technically oriented team.",
    "activity": [
      ("May 6", "HeySam", "Vista ERP export — MSHB/MSLB records, haul codes; vendor payables export first"),
      ("May 7", "Intercom", "Two open: approvals tab question + driver start-time notification"),
    ],
    "tickets": [],
    "risks": [],
    "systems": ["Vista ERP"],
  },
  {
    "name": "SILVERKING TRUCKING",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$55,200",
    "trucks": "70",
    "location": "Fort Myers, FL",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/18240007277",
    "what": "SW Florida aggregate hauling partner. Delivers aggregates for construction projects.",
    "connects_with": "Aggregate producers and construction sites in SW Florida",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Bulk Ticket Downloads (PDF)", "Hauler Pay Reporting", "Vendor Account Management"],
    "personality": "Passive engager — 4 open tickets but no meetings in 90 days. May be self-sufficient or quietly disengaging. No ARR recorded.",
    "activity": [],
    "tickets": [
      ("TRE-14126", "Bulk zip download silently hangs (lambda timeout)"),
      ("TRE-14118", "Bulk ticket PDF produces blank pages on large orders"),
      ("OPS-493", "Customers need to set vendors as managed accounts"),
      ("TRE-9998", "Redirected to Omni when downloading hauler pay report"),
    ],
    "risks": ["No ARR. No meetings in 90 days despite 4 open tickets."],
    "systems": [],
  },
  {
    "name": "TAPANI INC",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$39,900",
    "trucks": "50",
    "location": "Battle Ground, WA",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9410054012",
    "what": "Family-owned infrastructure construction in Pacific NW for 40+ years. Civil site, transportation, utilities.",
    "connects_with": "Hauler vendors and material suppliers; municipalities in WA/OR",
    "main_contacts": [("Riley Dettloff", "Driver — GPS issue contact")],
    "tread_features": ["Dispatch", "Payroll/Accounting Reporting", "HCSS Integration", "Vista ERP Integration", "GPS Tracking"],
    "personality": "Family-owned, meticulous. Multiple ERP integrations (HCSS, Vista, Fleetwatcher). Accounting + payroll accuracy are their core use case. PDF truncation frustration ongoing.",
    "activity": [("May 7", "HeySam", "Reporting issues — PDF truncation (1 of 77 rows cut off); ticket number column missing from dashboard report")],
    "tickets": [
      ("OPS-412", "Tapani Accounting Payroll Report (In Progress)"),
      ("TRE-14096", "Order #696191 job count mismatch"),
      ("TRE-9766", "System advanced driver to next load before leaving drop-off"),
      ("TRE-9493", "GPS signal drops until app reboot — Riley Dettloff"),
    ],
    "risks": [],
    "systems": ["HCSS", "Fleetwatcher", "Vista ERP"],
  },
  {
    "name": "TERRY EQUIPMENT COMPANY",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$25,000",
    "trucks": "—",
    "location": "Hillsboro, AL",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/18628365201",
    "what": "Family-owned hauling + contracting in Northern AL — material hauling, land clearing, site work.",
    "connects_with": "Material vendors and construction sites in Northern Alabama",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Budget Upload", "Training + Support Login"],
    "personality": "Family-owned, training-oriented. Recent session covered budget upload and support login. Day count discrepancy bug flagged.",
    "activity": [("May 1", "HeySam", "Training session — budget upload, support login, bug flagged, day count discrepancy")],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
  {
    "name": "THUNDERBOLT",
    "health": "gray",
    "customer_type": "Construction",
    "arr": "$38,364",
    "trucks": "30",
    "location": "Ottawa, ON",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409426578",
    "what": "Landscape construction in Ottawa area. 30 years experience. Hauls sand, gravel, topsoil, crushed stone.",
    "connects_with": "Sand, gravel, topsoil suppliers in Ottawa area; construction/landscape sites",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Project/Site Configuration", "Order Management"],
    "personality": "Recently onboarded. Pragmatic workaround user — appending letter suffixes to handle duplicate PO numbers. No CSM or ARR.",
    "activity": [("May 7", "HeySam", "Onboarding — project/site data config; duplicate PO number workaround (append letter suffixes)")],
    "tickets": [],
    "risks": ["No CSM. No ARR recorded."],
    "systems": [],
  },
  {
    "name": "TOP TIER TRUCKING",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$17,640",
    "trucks": "25",
    "location": "McKinney, TX",
    "csm": "unassigned",
    "owner": "unassigned",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/15989822262",
    "what": "Aggregate hauler in McKinney, TX. Planning shift toward owner-operator lease model. 25 trucks. 33 HubSpot contacts.",
    "connects_with": "Producers and construction companies in TX (McKinney area)",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Rate Tables", "Fuel Surcharge", "Settlement", "Owner-Operator Lease (planned)"],
    "personality": "Forward-thinking — planning owner-operator model shift. No CSM or owner assigned despite 33 contacts. Contract expired May 2025.",
    "activity": [("Apr 28", "HeySam", "Fuel surcharge tables + rate visibility confirmed; owner-operator lease rollout planned")],
    "tickets": [],
    "risks": ["No CSM or owner. Contract expired May 2025."],
    "systems": [],
  },
  {
    "name": "TRIO AGGREGATE HAULERS",
    "health": "green",
    "customer_type": "Hauler",
    "arr": "$25,200",
    "trucks": "70",
    "location": "—",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/48623222165",
    "what": "Aggregate hauling company, 70 trucks.",
    "connects_with": "Producers, quarries, construction sites",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Driver App"],
    "personality": "Quiet, healthy, no open issues. Low-touch account.",
    "activity": [],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
  {
    "name": "TWIN CITY HAULING",
    "health": "gray",
    "customer_type": "Hauler",
    "arr": "$5,500",
    "trucks": "15",
    "location": "Hastings, MN",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/45281990937",
    "what": "Transportation company in MN. 15 trucks.",
    "connects_with": "Material producers and shippers in Minnesota",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Driver App (struggling — drivers deleting it)"],
    "personality": "Driver app abandonment is a critical issue — drivers deleting app, one no-showed work. Two-ticket flow too complex. No follow-up in 4+ months after this was flagged.",
    "activity": [("Jan 9", "HeySam", "Critical: drivers deleting app, one no-showed work; two-ticket flow too complex")],
    "tickets": [],
    "risks": ["IMMEDIATE RISK: Driver app abandonment. No CSM. No follow-up in 4+ months."],
    "systems": [],
  },
  {
    "name": "WALKER AG GROUP",
    "health": "gray",
    "customer_type": "Agriculture",
    "arr": "$25,200",
    "trucks": "30",
    "location": "—",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/318330759903",
    "what": "Agricultural trucking company, 30 trucks. Newest account (Apr 2026). Go-live May 13.",
    "connects_with": "Agricultural producers + construction sites (Sullivan project is first active job)",
    "main_contacts": [("Sam Keck", "Controller — hands-on, flagged FSC gap")],
    "tread_features": ["Dispatch", "Vendor/Driver Onboarding", "Fuel Surcharge (gap — material subtotal not supported)"],
    "personality": "Brand-new — go-live May 13. Controller (Sam Keck) is hands-on and technically engaged. Fuel surcharge gap on material subtotal unresolved and could cause friction at launch.",
    "activity": [
      ("May 6", "HeySam", "Onboarding — project setup, dispatch, vendor/driver onboarding (Sullivan project, 30 trucks)"),
      ("May 8", "Gmail", "Sam Keck: FSC on material subtotal not supported; auto-scrape fuel index prices needed"),
    ],
    "tickets": [],
    "risks": ["CRITICAL ONBOARDING: Go-live May 13. Fuel surcharge gap (material subtotal) unresolved. No CSM."],
    "systems": [],
  },
  {
    "name": "WESTERN STATES CONTRACTING",
    "health": "green",
    "customer_type": "Construction",
    "arr": "$58,500",
    "trucks": "—",
    "location": "N. Las Vegas, NV",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/13371551301",
    "what": "Underground construction in NV — wet utilities, power, gas, grading. Enters 40–50 jobs/day manually.",
    "connects_with": "Haulers and material suppliers; municipalities in NV",
    "main_contacts": [("Crystal", "Finance — joining next review")],
    "tread_features": ["Dispatch", "Manual Order Entry (40-50/day)", "Phase Codes", "Reporting (Tons-per-hour)"],
    "personality": "High-volume manual entry is their core workflow. Bulk Excel import would be game-changing. Testing OptimoRoute as alternative. Crystal from finance joining next QBR.",
    "activity": [("Apr 28", "HeySam", "Excel bulk import request (OptimoRoute competitive risk) — no budget now; Crystal joining next review")],
    "tickets": [("TRE-14311", "Can't add new phase code from New Order form"), ("REP-116", "Tons-per-hour report: tons delivered / hours worked per project")],
    "risks": ["OptimoRoute is a competitive risk for bulk order entry."],
    "systems": ["OptimoRoute (testing)"],
  },
  {
    "name": "WILLIAMS TRUCKING CO.",
    "health": "red",
    "customer_type": "Hauler",
    "arr": "$70,000",
    "trucks": "50",
    "location": "Plant City, FL",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409426046",
    "what": "Transportation + roll-off services throughout Central FL. 50 trucks. 178 HubSpot contacts — largest in portfolio.",
    "connects_with": "Construction sites + producers in Central FL (hauler + roll-off services)",
    "main_contacts": [],
    "tread_features": ["Unknown — no recent activity"],
    "personality": "Red health. Sameer (CEO of Tread) is listed as CSM — not sustainable for day-to-day. Last contact Apr 9. Despite 178 contacts, nobody is engaging. $70K ARR at risk.",
    "activity": [],
    "tickets": [],
    "risks": ["RED HEALTH. CSM is Sameer (CEO). Last contact Apr 9. No recent activity despite $70K ARR."],
    "systems": [],
  },
  # ── New mid-market additions ──────────────────────────────────────────────
  {
    "name": "D CRUPI & SONS, INC.",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$12,797",
    "trucks": "—",
    "location": "Toronto, ON",
    "csm": "Latefa Redjouh",
    "owner": "Latefa Redjouh",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409467493",
    "what": "Family-owned Toronto construction company. Multiple GMs reflecting multi-division structure. Migrating to Horizon platform as of Jan 2026.",
    "connects_with": "Toronto construction sites",
    "main_contacts": [
      ("Dominic Crupi", "VP"),
      ("Dominic Passalacqua", "GM — primary contact, flagged missing start times"),
      ("Glenn Robillard", "GM"),
      ("Corrado Zotti", "Dispatcher"),
      ("Judith Niles", "Billing"),
    ],
    "tread_features": ["Dispatch Board", "Horizon (Migration)", "Ticket Reporting", "Vendor Management"],
    "personality": "Family-owned Toronto construction. CSM recently changed Adam Murray → Latefa Redjouh. Horizon migration underway. Issue: missing job start times in reports.",
    "activity": [("Jan 2026", "HeySam", "Horizon migration — 4 meetings. Domenic Passalacqua flagged missing start times in reports.")],
    "tickets": [],
    "risks": ["Horizon migration in progress — missing start times in reports is open pain point. CSM transition may cause continuity gap."],
    "systems": ["Horizon (migrating)"],
  },
  {
    "name": "GULFSHORE TRUCKING LLC",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$30,000",
    "trucks": "13",
    "location": "Fort Myers, FL",
    "csm": "unassigned",
    "owner": "Adam Murray",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409418652",
    "what": "Florida carrier-broker. 13 internal drivers + large subcontractor network. Charges 11% broker fee (10% for 'elite team'). Horizon migration planned.",
    "connects_with": "Florida subcontractors and carriers; mine/quarry customers requiring hauling authorization",
    "main_contacts": [
      ("Diana", "Daily platform ops — manages all day-to-day Tread activity"),
      ("Brandy Williams", "brandy@gulfshoretrucking.com"),
      ("Jaime Saavedra", "Dispatcher"),
      ("Amanda Evangelista", "Accounting Supervisor"),
    ],
    "tread_features": [
      "Dispatch Board (Horizon migration pending)",
      "Omni Reports (daily 7 PM + 10 AM dispatch reports)",
      "Ticket Upload / Management",
      "Driver App",
      "Rate Management (broker fee structure — 11% standard, 10% elite team)",
    ],
    "personality": "Unique carrier-broker model. KEY COMPLAINT: Tread increased their overhead — had to hire more staff to manage it. Heavy Omni report usage. 17 recorded meetings — very engaged but frustrated.",
    "activity": [
      ("Apr 2026", "HeySam", "Horizon migration discussion — Sameer attended. Migrating target list from legacy system."),
      ("Jan 2026", "HeySam", "Pre-launch Horizon check-in. Diana manages daily platform."),
    ],
    "tickets": [],
    "risks": ["Customer explicitly said Tread increased their overhead (hired more staff). Critical perception issue for renewal. Wants cleaner vendor rate options, audit trail for job removals, elite team column in reports."],
    "systems": ["Omni Reports"],
  },
  {
    "name": "R&R TRUCKING, INC.",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$78,718",
    "trucks": "—",
    "location": "Lewisville, TX",
    "csm": "unassigned",
    "owner": "Adam Murray",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/11583429760",
    "what": "Texas trucking with two divisions: pneumatics (9 job sites) and aggregates (1 job site). Mix of internal drivers and owner operators. Go-live Feb 12, 2026 on Horizon.",
    "connects_with": "Texas job sites (pneumatics + aggregate operations); owner operators",
    "main_contacts": [
      ("Alex", "Dispatcher — new to trucking dispatch (~4 months in)"),
      ("Michael Honey", "CFO — newly added to platform"),
    ],
    "tread_features": [
      "Horizon Dispatch",
      "Ticket Upload (per-load, real-time)",
      "Driver Mobile App (iPhone-dominant)",
      "Geofencing (Geo-Cert)",
      "Payroll Reports (commission-based per driver)",
      "Rate Management (internal vs. external driver visibility)",
    ],
    "personality": "Two-division TX trucking. Alex (dispatcher) is new to dispatch — ongoing support needed. Feature requests: driver groups for faster dispatching, prevent second ticket before first uploaded. Commission-based payroll.",
    "activity": [("Feb 12", "HeySam", "Go-live on Horizon. Payroll team still coming up to speed. Lead drivers trained in office.")],
    "tickets": [],
    "risks": ["New dispatcher (Alex) — experience gap creates support dependency. Payroll team training still needed."],
    "systems": ["Horizon"],
  },
  {
    "name": "UPPAL TRUCKING LTD",
    "health": "red",
    "customer_type": "Hauler",
    "arr": "$15,425",
    "trucks": "—",
    "location": "BC, Canada",
    "csm": "unassigned",
    "owner": "unassigned",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9282571541",
    "what": "Small owner-operator trucking company in BC, Canada. Minimal engagement data. Billing lapse (expired credit card).",
    "connects_with": "BC construction sites",
    "main_contacts": [],
    "tread_features": ["Basic Dispatch", "Driver App"],
    "personality": "Very limited engagement data. Credit card billing lapse suggests possible churn risk. No CSM assigned. Very little activity in HeySam.",
    "activity": [],
    "tickets": [],
    "risks": ["BILLING LAPSE — expired credit card. Minimal engagement. No CSM. Probable churn."],
    "systems": [],
  },
  {
    "name": "WERDCO BC INC.",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$179,214",
    "trucks": "—",
    "location": "Las Vegas, NV",
    "csm": "Latefa Redjouh",
    "owner": "Adam Murray",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9410030967",
    "what": "Las Vegas NV construction trucking. Mix of internal drivers and owner operators. QB Desktop integration in active weekly development since March 2026. Bills by ton; pays subhaulers by the hour.",
    "connects_with": "Las Vegas quarry/pit operators and construction sites; owner operators (leasers)",
    "main_contacts": [
      ("Brandon Conrad", "Operations Manager"),
      ("Jenn Pellegrino", "Accounts Payable"),
      ("Kole", "Primary Tread contact — day-to-day"),
      ("Tiffany", "User"),
    ],
    "tread_features": [
      "Dispatch Board (Tread Orders system)",
      "Ticket Upload",
      "GPS Tracking",
      "QuickBooks Desktop Integration (in active weekly development)",
      "Settlement Statements",
      "PO Management",
    ],
    "personality": "Bills by ton, pays subhaulers by hour — dual rate structure. Weekly QB Desktop integration calls since March 2026. Pain: duplicate orders lose driver assignments; afternoon system lag; unexpected invoice confused AP team.",
    "activity": [
      ("May 7", "HeySam", "Weekly QB Desktop integration call (most recent of 13 total meetings)"),
      ("Mar 2026", "HeySam", "CSM transition meeting. ConExpo invitation sent."),
    ],
    "tickets": [],
    "risks": ["QB Desktop integration still in weekly active development — incomplete. Afternoon system lag flagged. Unexpected invoice confused AP team."],
    "systems": ["QuickBooks Desktop (in progress)"],
  },
  {
    "name": "RAM-CO TRUCKING SERVICES",
    "health": "red",
    "customer_type": "Hauler",
    "arr": "$21K",
    "trucks": "—",
    "location": "Fort Lupton, CO",
    "csm": "unassigned",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409487761",
    "what": "CHURNED May 2026. Family-owned hauling company serving construction, landscaping, and oil field industries in the Denver area.",
    "connects_with": "",
    "main_contacts": [
      ("Denise Ramirez", "Co-Owner / Admin"),
      ("Juan Ramirez", "Operations"),
    ],
    "tread_features": [],
    "personality": "CHURNED May 2026. No CSM assigned at time of churn.",
    "activity": [],
    "tickets": [],
    "risks": ["CHURNED May 2026."],
    "systems": [],
  },
  {
    "name": "MIDTEX MATERIALS",
    "health": "red",
    "customer_type": "Producer",
    "arr": "$25.2K",
    "trucks": "—",
    "location": "Texas",
    "csm": "unassigned",
    "owner": "Will Amen",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/46773619375",
    "what": "CHURNED May 2026. Materials company based in Texas.",
    "connects_with": "",
    "main_contacts": [
      ("Mark Schmitt", "Manager"),
    ],
    "tread_features": [],
    "personality": "CHURNED May 2026. Previously managed under Adam Murray (CS deal).",
    "activity": [],
    "tickets": [],
    "risks": ["CHURNED May 2026."],
    "systems": [],
  },
  {
    "name": "TRANS-PHOS INC.",
    "health": "yellow",
    "customer_type": "Hauler",
    "arr": "$42,000",
    "trucks": "—",
    "location": "Florida",
    "csm": "Latefa Redjouh",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9248407095",
    "what": "Florida phosphate hauler. Large subcontractor network dispatched via AYR (main vendor broker). Multiple Florida quarry/mine locations. Digital ticketing implementation in active progress.",
    "connects_with": "AYR (main vendor broker); Florida phosphate mine and quarry sites; Monday.com for job automation",
    "main_contacts": [
      ("Kevin McGee", "kmcgee@transphos.com — Dispatcher / primary contact"),
    ],
    "tread_features": [
      "Digital Ticketing (implementation in progress)",
      "Scale + Digital Ticket Dual Workflow",
      "GPS Tracking",
      "Dispatch Board",
      "Monday.com API Integration (automating project/order creation)",
      "Great Plains (GP) Accounting Import",
      "AYR Vendor Portal",
    ],
    "personality": "Florida phosphate hauler with complex digital ticket rollout. 22 recorded meetings. Key pain: drivers reverting to paper when digital is live; duplicate tickets; GPS compliance. Monday.com automation for job creation is unique integration.",
    "activity": [
      ("Apr 29", "HeySam", "Most recent meeting — digital ticket rollout, GPS compliance, AYR vendor management"),
      ("May 2026", "Intercom", "Kevin McGee: digital ticket job stuck in 'planned' stage"),
      ("May 2026", "Intercom", "Spanish-language user: adding subcontractor truck question"),
    ],
    "tickets": [],
    "risks": ["Driver adoption of digital tickets — reverting to paper creates duplicate ticket problem. GPS location compliance issues (drivers not enabling location)."],
    "systems": ["Monday.com", "Great Plains (GP)", "AYR Vendor Portal"],
  },
  {
    "name": "WHITAKER TRANSPORTATION",
    "health": "green",
    "customer_type": "Construction",
    "arr": "$38,250",
    "trucks": "—",
    "location": "Oak Park, IL",
    "csm": "Latefa Redjouh",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/10235338519",
    "what": "Construction trucking company. HQ in Oak Park, IL with Utah field operations. Full fleet on driver app. Focus on non-billable time tracking and phase/cost codes.",
    "connects_with": "Utah construction sites; internal fleet and owner operators",
    "main_contacts": [],
    "tread_features": [
      "Dispatch Board",
      "GPS Tracking",
      "Non-Billable Time Tracking",
      "Water Truck Dispatch",
      "Phase / Cost Codes",
      "Driver Mobile App (full fleet adoption)",
    ],
    "personality": "Full driver buy-in on mobile app. Focus on operational efficiency and non-billable time tracking. Phase code / cost code setup is a key differentiator for them. Clean, healthy account.",
    "activity": [],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
]

# ── Enterprise company data (multi-location accounts) ──────────────────────
enterprise_companies = [
  {
    "name": "AMRIZE: SASK + WINNIPEG",
    "health": "gray",
    "customer_type": "Producer",
    "arr": "—",
    "trucks": "—",
    "location": "SK / MB, Canada",
    "csm": "Latefa Redjouh",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/53579309216",
    "what": "Part of Amrize (formerly Lafarge/Holcim). Aggregate/concrete producer. ~200–300 daily loads. Expansion to US/Canada regions under discussion.",
    "connects_with": "Haulers in Tread: Rock On Trucks, M&J Trucking, Shaw, Beach Transport (external hauler vendors)",
    "main_contacts": [("Sheila", "Fleet (reports to Randy)"), ("Randy", "Fleet Manager"), ("Tim", "On-site contact")],
    "tread_features": ["Dispatch", "Order Management", "Driver Notifications", "Internal + External Hauler Management"],
    "personality": "Enterprise mindset. Contract standardization across all Amrize regions is a priority — a good pilot here opens doors to more regions. Expansion-oriented.",
    "activity": [
      ("May 7", "HeySam", "Amrize Plano intro — ~200-300 daily loads; Tim on-site visit scheduled"),
      ("May 6", "HeySam", "Contract standardization across regions; Sheila now reports to Randy"),
    ],
    "tickets": [("TRE-13953", "Job notification times show wrong timezone for SK drivers")],
    "risks": ["Contract standardization across Amrize regions is key — this account is the pilot. Success here unlocks GTA, GVA, NCR expansion."],
    "systems": [],
  },
  {
    "name": "AMRIZE: NCR-TWIN CITIES",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$50,500",
    "trucks": "—",
    "location": "Minneapolis, MN",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/21225017493",
    "what": "Amrize (Holcim) Minneapolis / Fargo-Moorhead region. Dispatches to external partner fleets of 20-100 trucks each. Uses APEX scale/dispatch system. FleetWatcher contract resolution pending.",
    "connects_with": "External partner fleets in Minneapolis/Fargo-Moorhead; APEX system",
    "main_contacts": [
      ("Lucas Rath", "lucas.rath@amrize.com — Dispatcher"),
      ("Kristin Benallack", "kristin.benallack@amrize.com — User (also on SK account)"),
    ],
    "tread_features": ["Dispatch Board", "Vendor/Driver Management", "GPS Tracking", "APEX Integration"],
    "personality": "Part of Amrize enterprise rollout (Holcim spinoff). FleetWatcher contract resolution was gating full commitment. Very active with Tread team — 35+ meetings across all Amrize regions.",
    "activity": [
      ("May 2026", "Intercom", "Lucas Rath: unable to create new vendor/driver (conv 215474022198661)"),
      ("May 2026", "Intercom", "Kristin Benallack: open question (conv 215474197837742)"),
    ],
    "tickets": [],
    "risks": ["FleetWatcher contract resolution needed for full commitment. Onboarding still in progress."],
    "systems": ["APEX"],
  },
  {
    "name": "AMRIZE: GVA (BC)",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$85,545",
    "trucks": "—",
    "location": "Vancouver, BC",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9410106140",
    "what": "Amrize (Holcim) Greater Vancouver Area. BC quarry and construction aggregate operations. Part of Holcim/Amrize enterprise rollout. Uses APEX.",
    "connects_with": "BC quarry and construction partners; APEX system",
    "main_contacts": [],
    "tread_features": ["Dispatch Board", "Vendor/Driver Management", "GPS Tracking", "APEX Integration"],
    "personality": "Part of the broader Amrize multi-region enterprise deal. Active meetings across all Amrize regions. Holcim/Amrize rebrand transition in progress.",
    "activity": [],
    "tickets": [],
    "risks": ["Onboarding in progress — part of multi-region rollout. Success in other regions drives this account."],
    "systems": ["APEX"],
  },
  {
    "name": "AMRIZE: GTA",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$268,400",
    "trucks": "—",
    "location": "Toronto, ON",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409380785",
    "what": "Amrize (Holcim) Greater Toronto Area — largest Amrize contract ($268K ARR). Construction aggregate operations in the GTA. Uses APEX scale/dispatch system.",
    "connects_with": "GTA aggregate and construction partners; APEX system",
    "main_contacts": [],
    "tread_features": ["Dispatch Board", "Vendor/Driver Management", "GPS Tracking", "APEX Integration"],
    "personality": "Largest Amrize region and most revenue-significant. Part of multi-region enterprise rollout. 35+ meetings across Amrize accounts.",
    "activity": [],
    "tickets": [],
    "risks": ["Largest Amrize contract — any issues here have outsized impact on enterprise relationship."],
    "systems": ["APEX"],
  },
  {
    "name": "CEMEX USA",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$893,000",
    "trucks": "178+",
    "location": "Houston, TX",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409551643",
    "what": "Tread's largest account by ARR ($893K). Multinational cement company. Multiple Florida terminals + Houston HQ. 178+ drivers at one port site alone. Transitioning to Horizon platform.",
    "connects_with": "Florida cement terminals, Houston HQ; LoadMaster integration, APEX integration",
    "main_contacts": [
      ("Lionel Parenteau", "lionel.parenteau@cemex.com — active user, FL terminals"),
      ("Jose Montesino", "jose.montesino@cemex.com — active user"),
    ],
    "tread_features": [
      "Multi-Terminal Dispatch",
      "Horizon (Pilot — strategically critical migration)",
      "GPS Tracking",
      "Ticket Management",
      "LoadMaster Integration",
      "APEX Integration",
    ],
    "personality": "Flagship account — Tread's #1 by ARR. Multi-national complexity. Horizon migration is strategically critical for this account. Platform performance lag is actively flagged in Intercom — needs urgent attention.",
    "activity": [
      ("May 2026", "Intercom", "Lionel Parenteau: platform lagging severely — performance issues"),
      ("May 2026", "Intercom", "Jose Montesino: reporting issue"),
    ],
    "tickets": [],
    "risks": ["CRITICAL: Platform performance/lag flagged in Intercom for Tread's largest account. Horizon migration cannot stall. Any disruption to $893K ARR account is a company-level risk."],
    "systems": ["LoadMaster", "APEX", "Horizon (migrating)"],
  },
  {
    "name": "DUFFERIN AGGREGATES (CRH)",
    "health": "green",
    "customer_type": "Producer",
    "arr": "$89,800",
    "trucks": "—",
    "location": "Vaughan, ON",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9410030179",
    "what": "CRH subsidiary. Ontario aggregate operations with multiple quarry locations. 5,000+ employee parent company. Expansion deal closed.",
    "connects_with": "Ontario quarry and construction operations; CRH sister companies (incl. PJ Keating)",
    "main_contacts": [],
    "tread_features": ["Dispatch Board", "GPS Tracking", "Ticket Management", "Multi-Site Dispatch"],
    "personality": "CRH subsidiary (same parent as PJ Keating). Sophisticated aggregate enterprise. Established, healthy account with expansion deal closed. High-value relationship for potential CRH-wide growth.",
    "activity": [],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
  {
    "name": "NATIONAL LIME AND STONE",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$90,000",
    "trucks": "1,200/day",
    "location": "Findlay, OH",
    "csm": "unassigned",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409543587",
    "what": "Ohio aggregates company. Columbus quarry alone handles 1,200 trucks/day (all external haulers). Multiple Ohio quarry locations. Global Software (Wayscale/GDM) webhook integration in progress.",
    "connects_with": "External hauler fleets at Ohio quarry locations; Global Software / Wayscale scale system",
    "main_contacts": [],
    "tread_features": [
      "Scale Operator Workflows",
      "Ticket Matching",
      "Global Software / Wayscale Integration (webhook — in progress)",
      "Dispatch Board",
      "Multi-Quarry Dispatch",
    ],
    "personality": "Ohio aggregates with very high truck volume (1,200/day at Columbus). Complex webhook integration with Global Software still in progress. Sophisticated onboarding — one of Tread's highest-volume operations.",
    "activity": [],
    "tickets": [],
    "risks": ["Global Software / Wayscale integration complexity — webhook implementation not yet complete."],
    "systems": ["Global Software (Wayscale/GDM)"],
  },
  {
    "name": "TOMLINSON",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$244,481",
    "trucks": "—",
    "location": "Ottawa, ON",
    "csm": "Latefa Redjouh",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9282721287",
    "what": "R.W. Tomlinson Limited — Ottawa construction company. Multi-division operation. 180 HubSpot contacts. Active daily use by multiple dispatchers. Expansion deal of $50K closed.",
    "connects_with": "Ottawa construction sites; internal and external haulers",
    "main_contacts": [
      ("Shawn Baldwin", "sbaldwin@tomlinsongroup.com — Truck Dispatcher"),
      ("Robert Andeloro", "randeloro@tomlinsongroup.com — Dispatcher"),
      ("Gina Ashley", "gashley@tomlinsongroup.com — User"),
      ("truckdispatch", "truckdispatch@tomlinsongroup.com — Shared dispatch inbox"),
    ],
    "tread_features": ["Dispatch Board", "Bulk Edit", "Stagger Time", "Scale Operator Permissions", "Vendor Management"],
    "personality": "Long-standing Ottawa construction company. High-volume daily user. 4 simultaneous open Intercom issues — highest support burden in portfolio. Recurring dispatch-to-driver info mismatch. Engaged but frustrated.",
    "activity": [
      ("May 2026", "Intercom", "Shawn Baldwin: cannot create vendor (conv 215474227880490)"),
      ("May 2026", "Intercom", "Session load error — 'Something Went Wrong' (conv 215474223429030)"),
      ("May 2026", "Intercom", "Robert Andeloro: dispatch/driver info mismatch — recurring (conv 215474027908029)"),
      ("May 2026", "Intercom", "Civic address input bug on job lines (conv 215474022268366)"),
    ],
    "tickets": [],
    "risks": ["4 simultaneous open Intercom issues — most active support burden in portfolio. Recurring dispatch-to-driver info mismatch unresolved."],
    "systems": [],
  },
  {
    "name": "GRANITE CONSTRUCTION (SOCAL)",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$27,600",
    "trucks": "—",
    "location": "Indio, CA",
    "csm": "Latefa Redjouh",
    "owner": "unassigned",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/35011728253",
    "what": "SoCal division of Granite Construction (large public infrastructure company). High Desert job in progress.",
    "connects_with": "RHZ Trucking (hauler vendor set up in Tread); High Desert haulers",
    "main_contacts": [],
    "tread_features": ["Dispatch"],
    "personality": "Division of large public company with only 1 HubSpot contact. Low local engagement. Last meeting Feb 9 — outside 90-day window.",
    "activity": [],
    "tickets": [],
    "risks": ["No HubSpot owner. Last meeting Feb 9 (outside 90-day window). No recent contact."],
    "systems": [],
  },
  {
    "name": "RPM xCONSTRUCTION",
    "health": "yellow",
    "customer_type": "Construction",
    "arr": "$105,000",
    "trucks": "250",
    "location": "McKinney, TX",
    "csm": "Latefa Redjouh",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409426042",
    "what": "TX-based design + construction — excavation, road, retaining walls. Turnkey infrastructure.",
    "connects_with": "Quality Trucking (AR) — primary sub-hauler; Nickel Rock (300 trucks + GLS dispatch); Viewpoint ERP",
    "main_contacts": [("Daniel Byrne", "Director IT & Data Analytics"), ("Riley Todd", ""), ("Jordan Frasier", ""), ("Miguel Bautista", "Project Manager")],
    "tread_features": ["Dispatch", "Viewpoint API Integration", "Order Management", "Reporting", "Bulk Ticket"],
    "personality": "Technically demanding — IT Director (Daniel Byrne) drives integration. 15+ open feature requests. API site status bug causing live operational issues.",
    "activity": [("Apr 28", "HeySam", "Touchbase — Viewpoint API site status bug (overriding active/inactive); bulk ticket regression; 15 open requests")],
    "tickets": [],
    "risks": ["API site status bug is active operational risk. 15+ open feature/fix requests."],
    "systems": ["Viewpoint ERP"],
  },
  {
    "name": "STATEWIDE MATERIALS",
    "health": "yellow",
    "customer_type": "Mixed",
    "arr": "$204,160",
    "trucks": "350",
    "location": "Manor, TX",
    "csm": "Tim Chung",
    "owner": "unassigned",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/45942726542",
    "what": "Leading aggregate hauler in TX. Gravel, limestone, sand. Largest single-location fleet (350 trucks). 2nd-highest ARR ($204K).",
    "connects_with": "i-14 Aggregates (subcontractor for truck-sharing); construction companies across TX",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Sage Integration", "Order Management", "Vendor Management", "Reporting"],
    "personality": "Largest fleet, 2nd highest ARR. Data cleanup in progress (blocked until May 15). Sage export complete. Complex operations with subcontractor truck-sharing.",
    "activity": [
      ("May 7", "HeySam", "Weekly — data cleanup blocked until May 15; Sage export complete; test order planned"),
      ("May 8", "Gmail", "i-14 Aggregates subcontractor truck-sharing question — resolved"),
    ],
    "tickets": [("PRO-1071", "Statewide issues reported Apr 30 (sub-issues TRE-8524, PRO-1072)"), ("OPS-493", "Vendor managed accounts feature")],
    "risks": ["No HubSpot owner."],
    "systems": ["Sage"],
  },
  {
    "name": "TILCON CT INC",
    "health": "green",
    "customer_type": "Producer",
    "arr": "$95,000",
    "trucks": "175",
    "location": "New Britain, CT",
    "csm": "Tim Chung",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9327210246",
    "what": "Aggregates, asphalt, and concrete supplier/contractor in CT. 175 trucks. Apex ERP integration active. Largest contact roster in portfolio (83 HubSpot contacts).",
    "connects_with": "Hired haulers (LX drivers + third-party vendors); construction contractors in CT",
    "main_contacts": [],
    "tread_features": ["Dispatch", "Apex ERP Integration", "Wait-Time Reporting (Omni)", "Sub-Hauler (LX) Management", "Driver Time Approvals"],
    "personality": "Large, complex operation with 83 contacts. Apex integration active but has wrong-customer bug. LX sub-hauler visibility gap ongoing. Engaged with Tim Chung.",
    "activity": [("Apr 23", "HeySam", "Check-in — expected tons bug, Apex address sync bug, sub-hauler visibility gap for LX drivers")],
    "tickets": [
      ("TRE-13944", "Apex: Wrong customer pulled on several orders (In Progress)"),
      ("REP-107", "Tilcon wait time dashboard in Omni"),
      ("REP-16", "Duplicate rows in Driver Time Approvals Report"),
      ("TRE-9742", "Priority list: hourly trucking, phase codes, foreman approval"),
    ],
    "risks": [],
    "systems": ["Apex ERP"],
  },
  {
    "name": "UNITED STATES LIME & MINERALS",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$62,550",
    "trucks": "75",
    "location": "Dallas, TX",
    "csm": "Tim Chung",
    "owner": "Max Marhenke",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409426585",
    "what": "Public company. Lime + limestone products. Plants in AR, CO, LA, OK, TX. Complex Slurry/Vendor-Assigns-Driver GPS workflows.",
    "connects_with": "Haulers and industrial customers across multi-state operations; Motive telematics integration",
    "main_contacts": [("Peter Bailey", "txdispatch@uslm.com — customer filtering issue")],
    "tread_features": ["Dispatch", "GPS / Live Tracking", "Motive Integration", "Slurry + Vendor-Assigns-Driver Workflows", "Auto-Bind (in progress)"],
    "personality": "Public company with enterprise expectations. Complex multi-plant GPS workflows. DFW visibility still unresolved. Technical team engaged. Motive breadcrumb bug open.",
    "activity": [("Apr 22", "HeySam", "GPS tracking session — Slurry/Vendor-Assigns-Driver workflows; DFW visibility issue unresolved")],
    "tickets": [
      ("TRE-14322", "Activate auto-bind flag for DFW + Houston (In Progress)"),
      ("TRE-14392", "Motive telematics: trucks on live map but no breadcrumbs"),
      ("TRE-14032", "Customer filtering issue — Peter Bailey, txdispatch@uslm.com"),
    ],
    "risks": [],
    "systems": ["Motive"],
  },
  {
    "name": "VOLKER STEVIN CONTRACTING",
    "health": "gray",
    "customer_type": "Construction",
    "arr": "$250,200",
    "trucks": "—",
    "location": "Calgary, AB",
    "csm": "Tim Chung",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9410030737",
    "what": "Canadian highways + infrastructure contracting in AB. Highest-ARR single-location account ($182K). Microsoft 365 integration live.",
    "connects_with": "Vendor haulers: Burnco (and ~50-60 others being onboarded); government highway projects",
    "main_contacts": [("Tracy", "Dispatch training lead")],
    "tread_features": ["Dispatch", "Microsoft 365 Integration", "Vendor Onboarding (50-60 vendors)", "Workflow / Best Practice Training"],
    "personality": "High-value critical onboarding. Tracy driving dispatch training. Vendor invites to Burnco and ~50 others must go out immediately. On-site training at Sandman venue.",
    "activity": [
      ("May 8", "HeySam", "Vendor onboarding for ~50-60 vendors (Burnco etc.); venue scheduling risk at Sandman"),
      ("May 7", "HeySam", "Tracy: dispatch workflow training ahead of go-live"),
    ],
    "tickets": [],
    "risks": ["CRITICAL ONBOARDING: Vendor invites must go out now. Sandman venue scheduling risk."],
    "systems": ["Microsoft 365"],
  },
  {
    "name": "PJ KEATING CO",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$50,000",
    "trucks": "—",
    "location": "Lunenburg, MA",
    "csm": "Tim Chung",
    "owner": "Tim Chung",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/9409255817",
    "what": "CRH subsidiary. Asphalt plant in Lunenburg, MA. All external haulers — 30-35 trucks per job from 6-7 vendor companies being onboarded. Group Buy deal structure.",
    "connects_with": "6-7 vendor hauler fleets; APEX scale system; Viewpoint payroll",
    "main_contacts": [],
    "tread_features": [
      "Dispatch Board",
      "Vendor Onboarding (6-7 hauler fleets)",
      "APEX Integration",
      "Viewpoint Payroll Upload",
      "Ticket Management",
    ],
    "personality": "CRH subsidiary (same parent as Dufferin Aggregates). Early-stage Tread user as of spring 2026. All external haulers — no internal fleet. Vendor adoption is the make-or-break metric.",
    "activity": [("Mar-Apr 2026", "HeySam", "Just started dispatching. Onboarding 6-7 vendor fleets. APEX and Viewpoint integrations in setup.")],
    "tickets": [],
    "risks": ["All-external-hauler model — vendor adoption is critical. Any friction in vendor onboarding could stall the account."],
    "systems": ["APEX", "Viewpoint"],
  },
  {
    "name": "HOLCIM - NORTH CENTRAL (FARGO)",
    "health": "yellow",
    "customer_type": "Producer",
    "arr": "$50,500",
    "trucks": "—",
    "location": "Fargo, ND",
    "csm": "Latefa Redjouh",
    "owner": "unassigned",
    "hubspot": "",
    "what": "North Central Region (Fargo) division of Holcim, one of the world's largest building materials companies. Aggregate and cement producer.",
    "connects_with": "",
    "main_contacts": [],
    "tread_features": ["Dispatch", "GPS Tracking"],
    "personality": "",
    "activity": [],
    "tickets": [],
    "risks": [],
    "systems": [],
  },
  {
    "name": "ZEMBA INC.",
    "health": "red",
    "customer_type": "Construction",
    "arr": "$20,700",
    "trucks": "—",
    "location": "Zanesville, OH",
    "csm": "unassigned",
    "owner": "unassigned",
    "hubspot": "https://app.hubspot.com/contacts/21383822/record/0-2/15668643591",
    "what": "CHURNED April 2026. Ohio construction conglomerate with 15 business verticals (aggregate, concrete, trucking, water, sewer, septic). Left Tread for Command Cloud to consolidate all verticals.",
    "connects_with": "Ohio construction sites (historical)",
    "main_contacts": [],
    "tread_features": ["Dispatch Board (historical — now on Command Cloud)"],
    "personality": "CHURNED — April 2026. Migrated to Command Cloud for organizational consolidation across 15 verticals. Key gap: Tread lacked aggregate + scale integration. Relationship remains warm — open to reconnecting. Track for potential win-back.",
    "activity": [
      ("Apr 27", "HeySam", "Churn conversation — confirmed full move to Command Cloud. Had been running both systems for 30 days before deciding."),
    ],
    "tickets": [],
    "risks": ["CHURNED. Missing feature: aggregate + scale integration. Win-back opportunity if Tread adds this capability. Contact at CES / CONEX."],
    "systems": ["Command Cloud (current)"],
  },
]

# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, line=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if not line:
        shape.line.fill.background()
        shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_textbox(slide, l, t, w, h, text, fsize, bold, color, align, wrap, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fsize)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_hyperlink_to_run(run, url, para):
    rPr = run._r.get_or_add_rPr()
    rId = para.part.relate_to(
        url,
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
        is_external=True
    )
    hl = etree.SubElement(rPr, qn('a:hlinkClick'))
    hl.set(qn('r:id'), rId)


def add_link_tb(slide, l, t, w, h, display, url, fsize, color, align=PP_ALIGN.RIGHT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = display
    r.font.size = Pt(fsize)
    r.font.color.rgb = color
    r.font.underline = True
    add_hyperlink_to_run(r, url, p)
    return tb


def section_head(tf, label, first=False, fsize=13):
    if not first:
        sp = tf.add_paragraph()
        sp.text = ""
        sp.space_before = Pt(max(3, fsize - 8))
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label.upper()
    r.font.size = Pt(max(9, fsize - 1))
    r.font.bold = True
    r.font.color.rgb = NAVY
    return p


def slide_fsize(c):
    """Estimate content line count and return a font size that fits the slide."""
    feats       = FEATURE_OVERRIDES.get(c['name'], c.get('tread_features', []))
    support_ids = INTERCOM_SUPPORT.get(c['name'], [])

    def lines(text, cpl):   # chars-per-line estimate
        return max(1, (len(text) + cpl - 1) // cpl)

    # Right column
    r = 0
    r += 1 + min(len(feats), 6)
    if c.get('activity'):
        r += 1
        for d, s, sm in c['activity'][:3]:
            r += lines(f"[{d}] {s} — {sm}"[:140], 85)
    r += 1
    for tid, tdesc in c['tickets'][:5]:
        r += lines(f"{tid} — {tdesc}"[:100], 85)
    if support_ids:
        r += 1 + min(len(support_ids), 6)
    if c.get('systems'):
        r += 2

    # Left column
    l = 0
    l += 1 + lines(c['what'][:300], 58)
    if c.get('connects_with'):
        l += 1 + lines(c['connects_with'][:200], 58)
    contacts = c.get('main_contacts', [])
    if contacts:
        l += 1
        for nm, ti in contacts[:4]:
            l += lines(f"{nm} — {ti}"[:70], 58)

    total = max(r, l)
    if total <= 22: return 13
    if total <= 26: return 12
    if total <= 30: return 11
    if total <= 35: return 10
    return 9


def body_para(tf, text, color=DARK, italic=False, fsize=13):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fsize)
    r.font.color.rgb = color
    r.font.italic = italic
    return p


def bullet_para(tf, text, color=DARK, fsize=13):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "• " + text
    r.font.size = Pt(fsize)
    r.font.color.rgb = color
    return p


def initial_badge(slide, x, y, size, letter, bg_color):
    """Colored square with white initial letter — used as visual logo."""
    add_rect(slide, x, y, size, size, bg_color)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(size), Inches(size))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = letter
    r.font.size = Pt(int(size * 26))
    r.font.bold = True
    r.font.color.rgb = WHITE


# ── Build presentation ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── MAP SLIDE FUNCTION ─────────────────────────────────────────────────────
_MAP_CUSTOMERS = [
    # (name, lat, lon, health, segment)   segment: 'mm' or 'ent'
    # ── Mid-market ──────────────────────────────────────────────────────────
    ("4M TRUCKING",                   30.27, -97.74,  "yellow", "mm"),
    ("AHS",                           29.19, -82.13,  "yellow", "mm"),
    ("ARIZONA AGGREGATE SOLUTIONS",   33.45, -112.07, "yellow", "mm"),
    ("BRINKS LAND IMPROVEMENT",       43.96, -122.73, "green",  "mm"),
    ("BUESING CORP",                  33.60, -112.20, "green",  "mm"),
    ("CERUTTI & SONS TRANSPORTATION", 36.75, -119.77, "gray",   "mm"),
    ("CHARLES H CARTER & SON",        38.37, -88.36,  "green",  "mm"),
    ("DANIELA TRUCKING & GRADING",    27.80, -82.60,  "yellow", "mm"),
    ("DIAMOND MATERIALS",             39.74, -75.54,  "gray",   "mm"),
    ("EPIC MATERIALS INC",            28.70, -81.20,  "green",  "mm"),
    ("FLASH TRUCKING / GOLF AGRONOMICS", 26.76, -81.43, "green", "mm"),
    ("GEORGE J. IGEL & CO.",          39.96, -82.99,  "gray",   "mm"),
    ("GERNATT ASPHALT PRODUCTS",      42.48, -79.05,  "green",  "mm"),
    ("GRANITE CONSTRUCTION (SOCAL)",  33.72, -116.21, "yellow", "ent"),
    ("GULFSHORE TRUCKING LLC",        26.50, -81.90,  "yellow", "mm"),
    ("IROQUOIS BAR CORPORATION",      42.82, -78.83,  "green",  "mm"),
    ("LOBO LOGISTICS",                39.73, -104.83, "red",    "mm"),
    ("MANSTEEL REBAR LTD.",           43.87, -79.43,  "yellow", "mm"),
    ("HOLCIM - NORTH CENTRAL (FARGO)", 46.88, -96.78, "yellow", "ent"),
    ("MARCC TRUCKING",                32.78, -96.80,  "green",  "mm"),
    ("MMC MATERIALS INC",             32.41, -90.13,  "green",  "mm"),
    ("N.S. TRUCKING INC.",            28.39, -80.74,  "green",  "mm"),
    ("PETERSON COMPANIES",            45.37, -92.92,  "yellow", "mm"),
    ("PINERIDGE FARMS INC.",          21.31, -157.86, "yellow", "mm"),
    ("PRIME AGGREGATE TRANSPORTATION",31.00, -98.00,  "yellow", "mm"),
    ("QUALITY TRUCKING",              34.75, -92.29,  "green",  "mm"),
    ("RHINO TRUCKING INC.",           39.77, -86.16,  "green",  "mm"),
    ("ROCK ON TRUCKS",                45.55, -94.22,  "yellow", "mm"),
    ("RONYX LOGISTICS LLC",           41.88, -87.63,  "green",  "mm"),
    ("RPM xCONSTRUCTION",             33.20, -96.64,  "yellow", "ent"),
    ("R.W. DUNTEMAN CO.",             41.93, -87.99,  "gray",   "mm"),
    ("SILVERKING TRUCKING",           26.64, -81.87,  "green",  "mm"),
    ("STATEWIDE MATERIALS",           30.34, -97.55,  "yellow", "ent"),
    ("TAPANI INC",                    45.78, -122.52, "yellow", "mm"),
    ("TERRY EQUIPMENT COMPANY",       34.65, -86.97,  "yellow", "mm"),
    ("THUNDERBOLT",                   45.42, -75.70,  "gray",   "mm"),
    ("TILCON CT INC",                 41.66, -72.78,  "green",  "ent"),
    ("TOP TIER TRUCKING",             33.30, -96.75,  "green",  "mm"),
    ("TWIN CITY HAULING",             44.74, -92.86,  "gray",   "mm"),
    ("UNITED STATES LIME & MINERALS", 32.90, -97.10,  "yellow", "ent"),
    ("VOLKER STEVIN CONTRACTING",     51.05, -114.07, "gray",   "ent"),
    ("WESTERN STATES CONTRACTING",    36.20, -115.12, "green",  "mm"),
    ("WERDCO BC INC.",                36.17, -115.20, "yellow", "mm"),
    ("WILLIAMS TRUCKING CO.",         28.02, -82.13,  "red",    "mm"),
    ("D CRUJI & SONS, INC.",          43.65, -79.38,  "yellow", "mm"),
    ("PJ KEATING CO",                 42.60, -71.71,  "yellow", "ent"),
    ("R&R TRUCKING, INC.",            33.05, -97.00,  "yellow", "mm"),
    ("UPPAL TRUCKING LTD",            49.10, -122.85, "gray",   "mm"),
    # ── Enterprise ─────────────────────────────────────────────────────────
    ("AMRIZE: SASK + WINNIPEG",       50.50, -104.60, "gray",   "ent"),
    ("AMRIZE: NCR-TWIN CITIES",       44.98, -93.27,  "yellow", "ent"),
    ("AMRIZE: GVA (BC)",              49.25, -123.10, "yellow", "ent"),
    ("AMRIZE: GTA",                   43.72, -79.55,  "yellow", "ent"),
    ("CEMEX USA",                     29.76, -95.37,  "yellow", "ent"),
    ("DUFFERIN AGGREGATES (CRH)",     43.85, -79.52,  "green",  "ent"),
    ("NATIONAL LIME AND STONE",       41.04, -83.65,  "yellow", "ent"),
    ("TOMLINSON",                     45.52, -75.80,  "yellow", "ent"),
    ("TRANS-PHOS INC.",               27.90, -82.00,  "yellow", "mm"),
    ("WHITAKER TRANSPORTATION",       41.88, -87.50,  "green",  "mm"),
    ("ZEMBA INC.",                    39.94, -82.01,  "red",    "ent"),
]

_ENT_LABELS = {
    "AMRIZE: SASK + WINNIPEG":  "Amrize SK/MB",
    "AMRIZE: NCR-TWIN CITIES":  "Amrize Twin Cities",
    "AMRIZE: GVA (BC)":         "Amrize GVA",
    "AMRIZE: GTA":              "Amrize GTA",
    "CEMEX USA":                "CEMEX",
    "DUFFERIN AGGREGATES (CRH)":"Dufferin",
    "NATIONAL LIME AND STONE":  "Natl Lime",
    "TOMLINSON":                "Tomlinson",
    "ZEMBA INC.":               "Zemba",
    "GRANITE CONSTRUCTION (SOCAL)": "Granite Const.",
    "RPM xCONSTRUCTION": "RPM xConst.",
    "STATEWIDE MATERIALS": "Statewide",
    "TILCON CT INC": "Tilcon CT",
    "UNITED STATES LIME & MINERALS": "US Lime",
    "VOLKER STEVIN CONTRACTING": "Volker Stevin",
    "PJ KEATING CO": "PJ Keating",
    "HOLCIM - NORTH CENTRAL (FARGO)": "Holcim NCR",
}


def make_map_slide():
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.patches import Polygon as MplPolygon
    from matplotlib.collections import PatchCollection
    import numpy as np
    from collections import defaultdict

    ctx = ssl.create_default_context()
    ctx.check_hostname = False
    ctx.verify_mode = ssl.CERT_NONE

    def fetch_json(url):
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        return json.loads(urllib.request.urlopen(req, context=ctx, timeout=20).read())

    print("  Fetching geo data...")
    us_geo = fetch_json('https://raw.githubusercontent.com/PublicaMundi/MappingAPI/master/data/geojson/us-states.json')
    ca_geo = fetch_json('https://raw.githubusercontent.com/codeforamerica/click_that_hood/master/public/data/canada.geojson')

    DENSITY = {
        "Texas": 9, "Florida": 9,
        "Ontario": 6,
        "Illinois": 4, "Minnesota": 4,
        "Ohio": 3,
        "New York": 2, "California": 2, "Arizona": 2, "Nevada": 2,
        "British Columbia": 2,
        "Colorado": 1, "Delaware": 1, "Indiana": 1, "Oregon": 1,
        "Mississippi": 1, "Arkansas": 1, "Alabama": 1, "Connecticut": 1,
        "Washington": 1, "Massachusetts": 1, "Hawaii": 1,
        "Saskatchewan": 1, "Alberta": 1, "Manitoba": 1,
    }

    def heat_color(n):
        if n == 0:  return '#0D1B2A'
        if n == 1:  return '#1A3A55'
        if n <= 2:  return '#1B5276'
        if n <= 4:  return '#1A6FA0'
        if n <= 6:  return '#00838F'
        return '#C87B0A'

    # Tighter zoom: CONUS + southern Canada, room for title at top
    LON_MIN, LON_MAX = -126, -63
    LAT_MIN, LAT_MAX =  22,  58

    FIG_W, FIG_H = 13.33, 7.5
    fig, ax = plt.subplots(figsize=(FIG_W, FIG_H))
    fig.patch.set_facecolor('#0D1B2A')
    ax.set_facecolor('#0D1B2A')
    ax.set_xlim(LON_MIN, LON_MAX)
    ax.set_ylim(LAT_MIN, LAT_MAX)
    ax.set_aspect('auto')
    ax.axis('off')
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)

    def draw_geo_heat(geo):
        for feat in geo['features']:
            n     = DENSITY.get(feat.get('properties', {}).get('name', ''), 0)
            fill  = heat_color(n)
            geom  = feat['geometry']
            rings = ([geom['coordinates'][0]] if geom['type'] == 'Polygon'
                     else [p[0] for p in geom['coordinates']])
            patches = []
            for ring in rings:
                pts = np.array(ring)
                if pts[:,0].max() < LON_MIN or pts[:,0].min() > LON_MAX: continue
                if pts[:,1].max() < LAT_MIN or pts[:,1].min() > LAT_MAX: continue
                patches.append(MplPolygon(pts, closed=True))
            if patches:
                ax.add_collection(PatchCollection(
                    patches, facecolor=fill, edgecolor='#2A4468', linewidth=0.5, zorder=1))

    draw_geo_heat(us_geo)
    draw_geo_heat(ca_geo)

    HC = {'green': '#43A047', 'yellow': '#FFB300', 'red': '#E53935', 'gray': '#78909C'}

    # ── Cluster spread: 1.2-degree grid bins → ring layout ──────────────────
    CELL, LON_ASP = 1.2, 1.5
    cells, hi_buf = defaultdict(list), []
    for co in _MAP_CUSTOMERS:
        name, lat, lon, health, seg = co
        if lon < -140:
            hi_buf.append(co)
        else:
            cells[(round(lat / CELL), round(lon / CELL))].append(co)

    def ring_pos(n, cx, cy):
        r = 0.48 + 0.09 * max(0, n - 3)
        return [(cx + r * np.cos(2*np.pi*i/n - np.pi/2),
                 cy + r * LON_ASP * np.sin(2*np.pi*i/n - np.pi/2))
                for i in range(n)]

    plot_rows = []
    for key, group in cells.items():
        cx = np.mean([c[1] for c in group])
        cy = np.mean([c[2] for c in group])
        if len(group) == 1:
            co = group[0]
            plot_rows.append((co[0], co[3], co[4], co[2], co[1], cy, cx, False))
        else:
            for i, co in enumerate(sorted(group, key=lambda x: x[4])):
                plat, plon = ring_pos(len(group), cx, cy)[i]
                plot_rows.append((co[0], co[3], co[4], plon, plat, cy, cx, True))

    # ── HI inset (bottom-left) ───────────────────────────────────────────────
    HX, HY = -122.5, 27.5
    ax.add_patch(plt.Rectangle((HX-0.2, HY-0.2), 3.0, 2.8,
                                lw=0.5, edgecolor='#2A4468', facecolor='#1A3A55', zorder=2))
    ax.text(HX+1.3, HY+2.4, 'HI', color='#6A9ABE', fontsize=8,
            ha='center', fontweight='bold', zorder=3)
    for i, co in enumerate(hi_buf):
        col = HC.get(co[3], HC['gray'])
        ax.scatter(HX+1.3, HY+0.5+i*0.6, s=65, c=col, marker='o',
                   zorder=5, edgecolors='white', linewidths=0.4)

    # ── Spokes then dots ─────────────────────────────────────────────────────
    for name, health, seg, plon, plat, cy, cx, clustered in plot_rows:
        if clustered:
            ax.plot([cy, plon], [cx, plat],
                    color='#3A5A7A', lw=0.5, alpha=0.5, zorder=2)

    for name, health, seg, plon, plat, cy, cx, clustered in plot_rows:
        col = HC.get(health, HC['gray'])
        if seg == 'ent':
            ax.scatter(plon, plat, s=260, c=col, marker='*', zorder=6,
                       edgecolors='white', linewidths=0.8)
        else:
            ax.scatter(plon, plat, s=70, c=col, marker='o', zorder=5,
                       edgecolors='white', linewidths=0.45, alpha=0.95)

    # Enterprise labels: offset away from dot using annotate arrows
    for name, health, seg, plon, plat, cy, cx, clustered in plot_rows:
        if seg != 'ent': continue
        col   = HC.get(health, HC['gray'])
        label = _ENT_LABELS.get(name, name)
        # choose offset direction to avoid map edges
        dx = 22 if plon < (LON_MIN + LON_MAX) / 2 else -22
        dy = 14
        ax.annotate(label, xy=(plon, plat),
                    xytext=(dx, dy), textcoords='offset points',
                    fontsize=6.5, color='white', fontweight='bold', zorder=8,
                    arrowprops=dict(arrowstyle='->', color=col, lw=0.7,
                                    shrinkA=4, shrinkB=3),
                    bbox=dict(boxstyle='round,pad=0.2', fc='#0A1628',
                              ec=col, lw=0.7, alpha=0.92))

    # ── City callouts: annotate with arrows, offset from cluster ────────────
    # (xy = cluster centre, xytext = label anchor in data coords)
    # Each tuple: (label_lon, label_lat, arrow_lon, arrow_lat, text, col)
    city_callouts = [
        # Texas
        (-93.8, 33.8, -96.6, 33.0, "Dallas / McKinney", '#C87B0A'),
        (-99.5, 29.2, -97.6, 30.2, "Austin area",        '#C87B0A'),
        (-93.5, 28.5, -95.4, 29.8, "Houston",            '#C87B0A'),
        # Florida
        (-79.8, 28.4, -82.0, 27.9, "Tampa / Plant City", '#C87B0A'),
        (-79.5, 26.2, -81.8, 26.5, "Fort Myers",         '#C87B0A'),
        (-80.2, 29.8, -81.9, 29.2, "Central FL",         '#C87B0A'),
        # Ontario
        (-76.5, 44.8, -79.4, 43.7, "Toronto / GTA",      '#00838F'),
        (-73.5, 46.0, -75.7, 45.4, "Ottawa",             '#00838F'),
        # Illinois
        (-84.8, 42.6, -87.7, 41.9, "Chicago",            '#00838F'),
        # Minnesota
        (-90.6, 44.4, -93.3, 44.9, "Twin Cities",        '#00838F'),
        (-91.8, 46.4, -94.2, 45.6, "St. Cloud",          '#00838F'),
    ]
    for llon, llat, alon, alat, txt, col in city_callouts:
        ax.annotate(txt, xy=(alon, alat), xytext=(llon, llat),
                    fontsize=7.5, color='white', fontweight='bold',
                    ha='center', va='center', zorder=9,
                    arrowprops=dict(arrowstyle='->', color=col, lw=0.8,
                                    shrinkA=3, shrinkB=2),
                    bbox=dict(boxstyle='round,pad=0.25', fc=col,
                              ec='none', alpha=0.88))

    # ── Legend (bottom-right, 2 columns) ────────────────────────────────────
    health_h = [
        mpatches.Patch(color=HC['green'],  label='Healthy'),
        mpatches.Patch(color=HC['yellow'], label='Needs Attention'),
        mpatches.Patch(color=HC['red'],    label='At Risk'),
        mpatches.Patch(color=HC['gray'],   label='Inactive'),
    ]
    heat_h = [
        mpatches.Patch(color='#C87B0A', label='Hotspot (7+)'),
        mpatches.Patch(color='#00838F', label='4–6 customers'),
        mpatches.Patch(color='#1A6FA0', label='2–3 customers'),
        mpatches.Patch(color='#1A3A55', label='1 customer'),
    ]
    ent_h = ax.scatter([], [], s=220, c='white', marker='*', label='★ Enterprise')
    mm_h  = ax.scatter([], [], s=60,  c='white', marker='o', label='● Mid-Market')
    leg = ax.legend(handles=health_h + [ent_h, mm_h] + heat_h,
                    loc='lower right', facecolor='#0A1628', edgecolor='#2A4468',
                    labelcolor='white', fontsize=8, ncol=2,
                    bbox_to_anchor=(0.995, 0.01), columnspacing=0.9,
                    handlelength=1.2, framealpha=0.92)
    leg.get_frame().set_linewidth(0.7)

    # ── Title (fixed screen position, top of axes) ──────────────────────────
    n_ent = sum(1 for *_, s in _MAP_CUSTOMERS if s == 'ent')
    n_mm  = sum(1 for *_, s in _MAP_CUSTOMERS if s == 'mm')
    ax.text(0.5, 0.980, 'Customer Map — North America',
            transform=ax.transAxes, color='white', fontsize=15, fontweight='bold',
            ha='center', va='top', zorder=10,
            bbox=dict(boxstyle='round,pad=0.3', fc='#0A1628', ec='none', alpha=0.88))
    ax.text(0.5, 0.948,
            f'{n_ent} Enterprise (★)  ·  {n_mm} Mid-Market (●)  '
            f'·  Dot color = account health  ·  State fill = customer density',
            transform=ax.transAxes, color='#7FAFD4', fontsize=8,
            ha='center', va='top', zorder=10,
            bbox=dict(boxstyle='round,pad=0.2', fc='#0A1628', ec='none', alpha=0.80))

    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=165, bbox_inches='tight',
                facecolor='#0D1B2A', pad_inches=0)
    buf.seek(0)
    plt.close(fig)

    sm = prs.slides.add_slide(blank)
    sm.shapes.add_picture(buf, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    print("  Map slide done.")


# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
add_rect(s1, 0, 0, 13.33, 7.5, NAVY)
add_rect(s1, 0, 6.6, 13.33, 0.9, HEALTH['green'])

# Tagline strip
tb = s1.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.9))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
_total = len(midmarket_companies) + len(enterprise_companies)
r.text = (f"{_total} customers  ·  Enterprise (multi-location) & Mid-Market (single-location)  ·  "
          "Health, contacts, features, activity, tickets")
r.font.size = Pt(12)
r.font.color.rgb = WHITE

# Title
add_textbox(s1, 1.0, 2.1, 11.33, 1.0, "Tread Customer Reference Deck",
            42, True, WHITE, PP_ALIGN.CENTER, False)
add_textbox(s1, 1.0, 3.2, 11.33, 0.6, "Customer Success Onboarding Reference  ·  May 2026",
            22, False, WHITE, PP_ALIGN.CENTER, False)
add_textbox(s1, 1.5, 4.0, 10.33, 0.5, "Sources: HubSpot  ·  HeySam  ·  Gmail  ·  Linear  ·  Intercom",
            14, False, RGBColor(190, 205, 225), PP_ALIGN.CENTER, False)

# Type legend chips — centered on slide
TYPES = [('Hauler', TYPE_COLOR['Hauler']), ('Producer', TYPE_COLOR['Producer']),
         ('Construction', TYPE_COLOR['Construction']), ('Agriculture', TYPE_COLOR['Agriculture']),
         ('Mixed', TYPE_COLOR['Mixed'])]
_chip_widths = [len(label) * 0.13 + 0.35 for label, _ in TYPES]
_total_chip_w = sum(_chip_widths) + 0.45 * (len(TYPES) - 1)
lx = (13.33 - _total_chip_w) / 2
for (label, col), cw in zip(TYPES, _chip_widths):
    add_rect(s1, lx, 4.8, cw, 0.38, col)
    tb2 = s1.shapes.add_textbox(Inches(lx), Inches(4.8), Inches(cw), Inches(0.38))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    lx += cw + 0.45

# ── Section slide helper ────────────────────────────────────────────────────
def make_section_slide(label, subtitle, bar_color):
    ss = prs.slides.add_slide(blank)
    add_rect(ss, 0, 0, 13.33, 7.5, NAVY)
    add_rect(ss, 0, 5.9, 13.33, 1.6, bar_color)
    add_textbox(ss, 0.8, 1.8, 11.73, 1.2, label, 56, True, WHITE, PP_ALIGN.CENTER, False)
    add_textbox(ss, 0.8, 3.2, 11.73, 0.6, subtitle, 20, False,
                RGBColor(0xBB, 0xCE, 0xE6), PP_ALIGN.CENTER, False)

# ── Roster slide helper ─────────────────────────────────────────────────────
_SHORT = {
    "CANTON CONCRETE (DUPLICATE)":       "CANTON CONCRETE",
    "FLASH TRUCKING / GOLF AGRONOMICS":  "FLASH / GOLF AGRONOMICS",
    "DANIELA TRUCKING & GRADING":        "DANIELA TRUCKING",
    "UNITED STATES LIME & MINERALS":     "US LIME & MINERALS",
    "ARIZONA AGGREGATE SOLUTIONS":       "AZ AGGREGATE SOLUTIONS",
    "CERUTTI & SONS TRANSPORTATION":     "CERUTTI & SONS TRANSP.",
    "PRIME AGGREGATE TRANSPORTATION":    "PRIME AGGREGATE TRANSP.",
    "WESTERN STATES CONTRACTING":        "WESTERN STATES CONTR.",
    "BRINKS LAND IMPROVEMENT":           "BRINKS LAND IMPROVE.",
    "GRANITE CONSTRUCTION (SOCAL)":      "GRANITE CONST. (SOCAL)",
    "VOLKER STEVIN CONTRACTING":         "VOLKER STEVIN CONTR.",
    "DUFFERIN AGGREGATES (CRH)":         "DUFFERIN AGGREGATES",
    "NATIONAL LIME AND STONE":           "NATL LIME & STONE",
    "THE NATIONAL LIME AND STONE CO.":   "NATL LIME & STONE",
    "WHITAKER TRANSPORTATION":           "WHITAKER TRANSP.",
}

_HEALTH_LABELS = [('green','HEALTHY'),('yellow','WATCH'),('red','AT RISK'),('gray','UNKNOWN')]

def make_roster_slide(cos_list, title, show_arr=False):
    sr = prs.slides.add_slide(blank)
    add_rect(sr, 0, 0, 13.33, 7.5, WHITE)
    add_textbox(sr, 0.3, 0.1, 10, 0.5, title, 28, True, NAVY, PP_ALIGN.LEFT, False)

    hc_r = {k: sum(1 for c in cos_list if c['health'] == k) for k in HEALTH}
    cx = 0.3
    for k, lbl in _HEALTH_LABELS:
        chip_text = f"  {lbl}: {hc_r[k]}  "
        chip_w = len(chip_text) * 0.095 + 0.1
        add_rect(sr, cx, 0.68, chip_w, 0.26, HEALTH[k])
        tb_c = sr.shapes.add_textbox(Inches(cx), Inches(0.68), Inches(chip_w), Inches(0.26))
        p_c = tb_c.text_frame.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        r_c = p_c.add_run()
        r_c.text = chip_text.strip()
        r_c.font.size = Pt(13)
        r_c.font.bold = True
        r_c.font.color.rgb = WHITE
        cx += chip_w + 0.12

    add_rect(sr, 0.3, 1.04, 12.73, 0.015, MGRAY)

    COL_W_R   = 3.08
    COL_GAP_R = 0.13
    COL_X_R   = [0.3 + i * (COL_W_R + COL_GAP_R) for i in range(4)]
    HDR_Y_R   = 1.12
    HDR_H_R   = 0.3
    LIST_Y_R  = HDR_Y_R + HDR_H_R + 0.06
    LIST_H_R  = 7.5 - LIST_Y_R - 0.15

    for i, (health_key, label) in enumerate(_HEALTH_LABELS):
        cx = COL_X_R[i]
        h_col = HEALTH[health_key]
        count = hc_r[health_key]
        add_rect(sr, cx, HDR_Y_R, COL_W_R, HDR_H_R, h_col)
        tb_hdr = sr.shapes.add_textbox(Inches(cx), Inches(HDR_Y_R), Inches(COL_W_R), Inches(HDR_H_R))
        p_hdr = tb_hdr.text_frame.paragraphs[0]
        p_hdr.alignment = PP_ALIGN.CENTER
        r_hdr = p_hdr.add_run()
        r_hdr.text = f"{label}  ({count})"
        r_hdr.font.size = Pt(14)
        r_hdr.font.bold = True
        r_hdr.font.color.rgb = WHITE

        group_cos = sorted([c for c in cos_list if c['health'] == health_key],
                           key=lambda x: x['name'])
        tb_lst = sr.shapes.add_textbox(Inches(cx + 0.08), Inches(LIST_Y_R),
                                       Inches(COL_W_R - 0.1), Inches(LIST_H_R))
        tf_lst = tb_lst.text_frame
        tf_lst.word_wrap = False
        first = True
        for co in group_cos:
            p = tf_lst.paragraphs[0] if first else tf_lst.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            display = _SHORT.get(co['name'], co['name'])
            if show_arr and co.get('arr') and co['arr'] != '—':
                display += f"  {co['arr']}"
            r.text = display
            r.font.size = Pt(12)
            r.font.color.rgb = DARK

# ── SLIDE 2: Customer Map ──────────────────────────────────────────────────
make_map_slide()

# ── SLIDE 3: Mid-Market Portfolio at a Glance ──────────────────────────────
_all_cos = midmarket_companies + enterprise_companies
make_roster_slide(midmarket_companies, "Mid-Market Accounts — Portfolio at a Glance")

# ── SLIDE 3: Watch List ─────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
add_rect(s3, 0, 0, 13.33, 7.5, WHITE)
add_rect(s3, 0, 0, 0.18, 7.5, HEALTH['red'])
add_textbox(s3, 0.42, 0.14, 12, 0.6, "Key Flags & Watch List", 30, True, DARK, PP_ALIGN.LEFT, False)

flags = [
    ("IMMEDIATE RISK — TRUST BREAKDOWN", HEALTH['red'],
     "Iroquois Bar Corp (dispatch issue open 37+ days)   |   Rock On Trucks (CEO trust breakdown, Sameer involved)   |   JW Golding (driver payout delayed Sept 30)   |   Twin City Hauling (app abandonment, 4mo no follow-up)"),
    ("RED HEALTH — SILENT CHURN", HEALTH['red'],
     "Lobo Logistics ($75.6K ARR — zero recent contact)   |   Williams Trucking ($70K ARR — CEO as CSM, last contact Apr 9)"),
    ("PLATFORM COMMITMENT AT RISK", WARN_T,
     "Marex (actively evaluating Turvo)   |   4M Trucking (COO frustrated, URGENT bug open, data integrity)   |   Pineridge Farms (QB sync has no timeline)   |   Peterson (paying for unused software)"),
    ("CRITICAL ONBOARDINGS THIS WEEK", WARN_T,
     "Volker Stevin (vendor invites to 50-60 vendors — Burnco etc. — needed now)   |   George Igel (pilot May 11)   |   Walker Ag (go-live May 13, fuel surcharge gap unresolved)"),
    ("OPEN INTERCOM ISSUES", TYPE_COLOR['Hauler'],
     "Rock On Trucks (4 open)   |   Ronyx Logistics (driver assign + 206-reply thread)   |   Diamond Materials (notification — open)   |   Iroquois Bar (37+ days)   |   MMC Materials (GPS)   |   R.W. Dunteman (2 open)"),
]
fy = 0.85
for heading, hcol, body in flags:
    add_textbox(s3, 0.42, fy, 12.7, 0.35, heading, 11, True, hcol, PP_ALIGN.LEFT, False)
    fy += 0.33
    add_textbox(s3, 0.62, fy, 12.5, 0.5, body, 13, False, DARK, PP_ALIGN.LEFT, True)
    fy += 0.62

# ── Per-company slides ──────────────────────────────────────────────────────

def build_slide(prs, c):
    fsize = slide_fsize(c)
    slide = prs.slides.add_slide(blank)
    h_col = HEALTH[c['health']]
    t_col = TYPE_COLOR.get(c['customer_type'], MED)

    # ─ Header bar ─
    HDR_H = 1.35
    add_rect(slide, 0, 0, 13.33, HDR_H, h_col)

    # Logo or initial badge
    BADGE = 1.1
    badge_x, badge_y = 0.12, 0.12
    logo_path = LOGO_PATHS.get(c['name'])
    add_rect(slide, badge_x, badge_y, BADGE, BADGE, WHITE)
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(badge_x), Inches(badge_y), Inches(BADGE), Inches(BADGE))
        except Exception:
            logo_path = None
    if not (logo_path and os.path.exists(logo_path)):
        letter = c['name'][0]
        tb_init = slide.shapes.add_textbox(Inches(badge_x), Inches(badge_y), Inches(BADGE), Inches(BADGE))
        p_i = tb_init.text_frame.paragraphs[0]
        p_i.alignment = PP_ALIGN.CENTER
        r_i = p_i.add_run()
        r_i.text = letter
        r_i.font.size = Pt(40)
        r_i.font.bold = True
        r_i.font.color.rgb = t_col

    # Company name
    tb_nm = slide.shapes.add_textbox(Inches(1.38), Inches(0.08), Inches(9.0), Inches(0.72))
    p_nm = tb_nm.text_frame.paragraphs[0]
    p_nm.alignment = PP_ALIGN.LEFT
    r_nm = p_nm.add_run()
    r_nm.text = c['name'][:52]
    r_nm.font.size = Pt(32)
    r_nm.font.bold = True
    r_nm.font.color.rgb = WHITE

    # Health + Type badges (right side)
    tb_ht = slide.shapes.add_textbox(Inches(10.0), Inches(0.06), Inches(3.2), Inches(0.55))
    p_ht = tb_ht.text_frame.paragraphs[0]
    p_ht.alignment = PP_ALIGN.RIGHT
    r_ht = p_ht.add_run()
    r_ht.text = c.get('customer_type', '').upper()
    r_ht.font.size = Pt(14)
    r_ht.font.bold = True
    r_ht.font.color.rgb = WHITE

    tb_h2 = slide.shapes.add_textbox(Inches(10.0), Inches(0.68), Inches(3.2), Inches(0.55))
    p_h2 = tb_h2.text_frame.paragraphs[0]
    p_h2.alignment = PP_ALIGN.RIGHT
    r_h2 = p_h2.add_run()
    r_h2.text = "● " + c['health'].upper()
    r_h2.font.size = Pt(14)
    r_h2.font.color.rgb = WHITE

    # ─ Info strip ─
    INFO_Y = HDR_H
    INFO_H = 0.55
    add_rect(slide, 0, INFO_Y, 13.33, INFO_H, LGRAY)
    parts = []
    if c.get('location') and c['location'] not in ('—', ''):
        parts.append(f"  {c['location']}")
    if c.get('arr') and c['arr'] != '—':
        parts.append(f"ARR: {c['arr']}")
    if c.get('trucks') and c['trucks'] != '—':
        parts.append(f"Trucks: {c['trucks']}")
    if c.get('csm'):
        parts.append(f"CSM: {c['csm']}")
    tenure = TENURE_APPROX.get(c['name'])
    if tenure:
        parts.append(f"Tenure: {tenure}")
    if c.get('go_live'):
        parts.append(f"Go-Live: {c['go_live']}")
    if c.get('next_qbr'):
        parts.append(f"Next QBR: {c['next_qbr']}")
    info_str = "   |   ".join(parts)

    tb_info = slide.shapes.add_textbox(Inches(0.2), Inches(INFO_Y), Inches(8.4), Inches(INFO_H))
    p_info = tb_info.text_frame.paragraphs[0]
    p_info.alignment = PP_ALIGN.LEFT
    r_info = p_info.add_run()
    r_info.text = info_str[:145]
    r_info.font.size = Pt(13)
    r_info.font.color.rgb = MED

    # Status badge (colored)
    _status = USAGE_STATUS.get(c['name'], '')
    _status_colors = {
        "Primary system": RGBColor(0x2E, 0x7D, 0x32),   # dark green
        "Onboarding":     RGBColor(0x15, 0x65, 0xC0),   # blue
        "Sporadic":       RGBColor(0xE6, 0x5C, 0x00),   # orange
        "Disengaged":     RGBColor(0xC6, 0x28, 0x28),   # red
    }
    if _status:
        _sc = _status_colors.get(_status, MED)
        tb_st = slide.shapes.add_textbox(Inches(8.7), Inches(INFO_Y + 0.05), Inches(2.1), Inches(0.45))
        p_st = tb_st.text_frame.paragraphs[0]
        p_st.alignment = PP_ALIGN.CENTER
        r_st = p_st.add_run()
        r_st.text = f"▶  {_status}"
        r_st.font.size = Pt(13)
        r_st.font.bold = True
        r_st.font.color.rgb = _sc

    if c.get('hubspot'):
        add_link_tb(slide, 10.95, INFO_Y + 0.05, 2.0, 0.45, "HubSpot →", c['hubspot'], 12, LINK)

    # Divider
    add_rect(slide, 0, INFO_Y + INFO_H, 13.33, 0.015, MGRAY)

    # ─ Content area ─
    CT = INFO_Y + INFO_H + 0.06   # content top
    CB = 6.15                      # content bottom
    CH = CB - CT

    # Left panel 0.2 → 5.7 (5.5" wide)
    L_W = 5.5
    tb_L = slide.shapes.add_textbox(Inches(0.2), Inches(CT), Inches(L_W), Inches(CH))
    tf_L = tb_L.text_frame
    tf_L.word_wrap = True

    # ABOUT
    section_head(tf_L, "About", first=True, fsize=fsize)
    body_para(tf_L, c['what'][:300], fsize=fsize)

    # CONNECTS WITH IN TREAD
    if c.get('connects_with'):
        section_head(tf_L, "Connects With in Tread", fsize=fsize)
        body_para(tf_L, c['connects_with'][:200], color=MED, italic=True, fsize=fsize)

    # MAIN CONTACTS
    contacts = c.get('main_contacts', [])
    if contacts:
        section_head(tf_L, "Main Contacts", fsize=fsize)
        for name, title in contacts[:4]:
            line = name if not title else f"{name}  —  {title}"
            bullet_para(tf_L, line[:70], fsize=fsize)

    # Vertical divider
    add_rect(slide, 5.78, CT, 0.015, CH, MGRAY)

    # Right panel 5.82 → 13.1 (7.28" wide)
    R_X = 5.82
    R_W = 7.28
    tb_R = slide.shapes.add_textbox(Inches(R_X), Inches(CT), Inches(R_W), Inches(CH))
    tf_R = tb_R.text_frame
    tf_R.word_wrap = True

    # TREAD FEATURES
    section_head(tf_R, "Key Tread Features Used", first=True, fsize=fsize)
    feats = FEATURE_OVERRIDES.get(c['name'], c.get('tread_features', []))
    for f in feats[:6]:
        bullet_para(tf_R, f[:80], fsize=fsize)
    if not feats:
        body_para(tf_R, "Unknown", color=MED, italic=True, fsize=fsize)

    # RECENT ACTIVITY
    if c.get('activity'):
        section_head(tf_R, "Recent Activity", fsize=fsize)
        for date, src, summary in c['activity'][:3]:
            p_act = tf_R.add_paragraph()
            p_act.alignment = PP_ALIGN.LEFT
            r_act = p_act.add_run()
            r_act.text = f"• [{date}] {src} — {summary}"[:140]
            r_act.font.size = Pt(fsize)
            r_act.font.color.rgb = DARK

    # OPEN TICKETS
    section_head(tf_R, "Open Tickets", fsize=fsize)
    if c['tickets']:
        for tid, tdesc in c['tickets'][:5]:
            tp = tf_R.add_paragraph()
            tp.alignment = PP_ALIGN.LEFT
            url = f"https://linear.app/treadapp/issue/{tid}"
            r_id = tp.add_run()
            r_id.text = tid
            r_id.font.size = Pt(fsize)
            r_id.font.color.rgb = LINK
            r_id.font.underline = True
            add_hyperlink_to_run(r_id, url, tp)
            r_desc = tp.add_run()
            r_desc.text = f" — {tdesc}"[:100]
            r_desc.font.size = Pt(fsize)
            r_desc.font.color.rgb = DARK
    else:
        body_para(tf_R, "None open", color=MED, italic=True, fsize=fsize)

    # OPEN SUPPORT CONVERSATIONS (INTERCOM — last 30 days)
    support_ids = INTERCOM_SUPPORT.get(c['name'], [])
    if support_ids:
        section_head(tf_R, "Open Support Tickets (Intercom — Last 30 Days)", fsize=fsize)
        cap = 6
        for conv_id in support_ids[:cap]:
            conv_url = f"https://app.intercom.com/a/apps/m48souwv/conversations/{conv_id}"
            sp = tf_R.add_paragraph()
            sp.alignment = PP_ALIGN.LEFT
            r_s = sp.add_run()
            r_s.text = f"#{conv_id[-8:]}"
            r_s.font.size = Pt(fsize)
            r_s.font.color.rgb = LINK
            r_s.font.underline = True
            add_hyperlink_to_run(r_s, conv_url, sp)
        if len(support_ids) > cap:
            body_para(tf_R, f"+ {len(support_ids)-cap} more open", color=MED,
                      italic=True, fsize=max(9, fsize - 2))

    # PRIOR/INTEGRATED SYSTEMS
    if c.get('systems'):
        section_head(tf_R, "Prior / Integrated Systems", fsize=fsize)
        body_para(tf_R, ", ".join(c['systems']), fsize=fsize)

    # ─ Footer ─
    FT_Y = 6.15
    FT_H = 1.35
    personality = c.get('personality', '')
    risks = c.get('risks', [])

    if risks:
        add_rect(slide, 0, FT_Y, 13.33, FT_H, WARN)
        tb_f = slide.shapes.add_textbox(Inches(0.25), Inches(FT_Y), Inches(12.83), Inches(FT_H))
        tf_f = tb_f.text_frame
        tf_f.word_wrap = True
        p_r = tf_f.paragraphs[0]
        p_r.alignment = PP_ALIGN.LEFT
        r_r = p_r.add_run()
        r_r.text = "RISK:  " + "   |   ".join(risks)[:240]
        r_r.font.size = Pt(max(10, fsize - 1))
        r_r.font.bold = True
        r_r.font.color.rgb = WARN_T
        if personality:
            p_p = tf_f.add_paragraph()
            r_p = p_p.add_run()
            r_p.text = "ENGAGEMENT: " + personality[:200]
            r_p.font.size = Pt(max(9, fsize - 2))
            r_p.font.italic = True
            r_p.font.color.rgb = WARN_T
    else:
        add_rect(slide, 0, FT_Y, 13.33, FT_H, LGRAY)
        tb_f = slide.shapes.add_textbox(Inches(0.25), Inches(FT_Y), Inches(12.83), Inches(FT_H))
        tf_f = tb_f.text_frame
        tf_f.word_wrap = True
        p_e = tf_f.paragraphs[0]
        p_e.alignment = PP_ALIGN.LEFT
        r_e = p_e.add_run()
        r_e.text = "ENGAGEMENT STYLE:  " + (personality or "Healthy, low-risk account.")[:220]
        r_e.font.size = Pt(max(10, fsize))
        r_e.font.italic = True
        r_e.font.color.rgb = MED

    return slide


# ── Section: Enterprise ────────────────────────────────────────────────────
if enterprise_companies:
    make_section_slide(
        "ENTERPRISE",
        f"{len(enterprise_companies)} accounts  ·  Multi-location / multi-division",
        HEALTH['green']
    )
    make_roster_slide(enterprise_companies,
                      "Enterprise Accounts — Portfolio at a Glance",
                      show_arr=True)
    for co in sorted(enterprise_companies, key=lambda x: x['name']):
        build_slide(prs, co)

# ── Section: Mid-Market ────────────────────────────────────────────────────
make_section_slide(
    "MID-MARKET",
    f"{len(midmarket_companies)} accounts  ·  Single-location",
    RGBColor(0x15, 0x65, 0xC0)
)
for co in sorted(midmarket_companies, key=lambda x: x['name']):
    build_slide(prs, co)

out = "/tmp/tread_customers.pptx"
prs.save(out)
import os
print(f"Done! {os.path.getsize(out)//1024} KB, {len(prs.slides)} slides")
